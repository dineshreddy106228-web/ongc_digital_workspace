from datetime import date
from decimal import Decimal
from types import SimpleNamespace

from app.core.services.inventory_monitoring import material_code, normalize_name


def test_material_code_preserves_or_restores_group_09_leading_zeroes():
    assert material_code("090001043") == "090001043"
    assert material_code(90001043) == "090001043"
    assert material_code("90001043.0") == "090001043"


def test_normalize_name_is_case_and_whitespace_insensitive():
    assert normalize_name("  Assam   Asset ") == "assam asset"
    assert normalize_name("ASSAM ASSET") == "assam asset"


def test_decimal_parser_handles_empty_and_formatted_numbers():
    from app.core.services.inventory_monitoring import _decimal

    assert _decimal("1,234.50") == Decimal("1234.50")
    assert _decimal("") is None


def _register(key, label, count, months, value_crore):  # noqa: D401
    rows = [
        {
            "group": "10", "code": f"1000{index:05d}", "description": "BARYTES API GRADE SPECIFICATION",
            "centre": "Sivasagar ST", "zone": "Assam", "qty": Decimal("40"), "uom": "MT",
            "value": Decimal(str(value_crore)) * Decimal("10000000"), "open_po": None, "open_pr": None,
            "months": None if months is None else Decimal(str(months)), "details": None,
        }
        for index in range(count)
    ]
    for row in rows:
        row["section"] = "Drilling Fluid Chemicals"
    return {
        "key": key, "label": label, "description": f"{label} register.", "count": count,
        "value": sum((row["value"] for row in rows), Decimal("0")), "rows": rows, "omitted": 0,
        "groups": [{"key": "DFC", "label": "Drilling Fluid Chemicals", "rows": rows, "total": count, "omitted": 0}],
    }


def _management_review_stub():
    crore = Decimal("10000000")
    return {
        "reporting_date": date(2026, 7, 31), "previous_date": date(2026, 6, 30), "row_limit": 250, "group_limit": 40,
        "thresholds": {"critical_low_stock_months": Decimal("1"), "low_stock_months": Decimal("3"), "slow_moving_months": Decimal("6"), "excess_stock_months": Decimal("12")},
        "kpis": {
            "total_value": 870 * crore, "prev_total_value": 810 * crore, "value_at_risk": 310 * crore,
            "prev_value_at_risk": 280 * crore, "at_risk_share": 35.6, "stockout_value": 42 * crore,
            "prev_stockout_value": 51 * crore, "stockout_share": 4.8,
            "centre_count": 46, "material_count": 812, "record_count": 4210, "top5_share": 61.2,
        },
        "health_mix": [{"key": "excess_stock", "label": "Excess", "value": 210 * crore, "count": 700, "share": 24.1}],
        "zones": [{"zone": "Assam", "value": 260 * crore, "share": 29.9, "prev": 250 * crore}],
        "centres_ranked": [{"name": f"Work Centre {index}", "zone": "Assam", "value": (90 - index) * crore, "share": 5.0} for index in range(20)],
        "movers": {"up": [{"name": "Ankleshwar DFS", "delta": 22 * crore}], "down": [{"name": "Mehsana DFS", "delta": -14 * crore}]},
        "entrants": [{"name": "Karaikal ST", "zone": "Southern", "value": 6 * crore}],
        "exits": [{"name": "Agartala DFS", "zone": "North Eastern", "value": 4 * crore}],
        "comparison": {
            "previous_date": date(2026, 6, 30), "gap_days": 31, "common_centres": 44,
            "entrant_count": 1, "entrant_value": 6 * crore, "exit_count": 1, "exit_value": 4 * crore,
            "like_for_like": 864 * crore, "like_for_like_prev": 806 * crore,
            "like_for_like_delta": 58 * crore, "like_for_like_change": 7.2, "is_default": True,
        },
        "high_value_materials": [
            {"code": "090001043", "description": "OIL WELL CEMENT CLASS G", "group": "09", "value": 60 * crore,
             "share": 6.9, "centres": 4, "months_low": Decimal("0.4"), "months_high": Decimal("18.2"),
             "section": "Drilling Fluid Chemicals"}
        ],
        "high_value_total": 60 * crore,
        "coverage_registers": [
            _register("critical_low_stock", "Critical low stock", 18, 0.4, 0.8),
            _register("excess_stock", "Excess stock", 20, 26, 2.2),
            _register("open_supply_with_high_stock", "Open PO / PR against high stock", 0, 14, 0),
        ],
        "supporting_registers": [_register("non_moving", "Non-moving materials", 33, None, 1.9)],
    }


def test_management_presentation_covers_every_register(monkeypatch, tmp_path):
    from pptx import Presentation

    import app.core.services.inventory_monitoring as monitoring
    from app.core.services.inventory_presentation import build_management_review_presentation

    monkeypatch.setattr(monitoring, "management_review_data", lambda reporting_date=None, compare_date=None, centre_ids=None: _management_review_stub())
    output, filename = build_management_review_presentation(str(tmp_path))

    assert filename == "ONGC Inventory Management Review 31 Jul 2026.pptx"
    titles = [
        shape.text_frame.text
        for slide in Presentation(output).slides
        for shape in slide.shapes
        if shape.has_text_frame and shape.text_frame.text
    ]
    assert "Materials above ₹ 1 Cr of inventory value" in titles
    assert "Not in Corporate Specification List" not in titles  # the stub has no unspecified material
    assert "Non-moving materials (1–13 of 33)" in titles
    assert "Critical low stock (1–13 of 18)" in titles
    assert "Excess stock (1–13 of 20)" in titles
    assert "Open PO / PR against high stock" in titles
    assert "Decisions sought from the review" in titles
    assert "Movement against the comparison period" in titles
    assert any("44 work centres reported stock in both" in title for title in titles)


def test_management_presentation_requires_a_published_period(monkeypatch, tmp_path):
    import pytest

    import app.core.services.inventory_monitoring as monitoring
    from app.core.services.inventory_presentation import build_management_review_presentation

    monkeypatch.setattr(monitoring, "management_review_data", lambda reporting_date=None, compare_date=None, centre_ids=None: {"reporting_date": None, "kpis": None})
    with pytest.raises(ValueError):
        build_management_review_presentation(str(tmp_path))


def _centre_record(code, description, group, months, value_crore, qty):
    return SimpleNamespace(
        material_code=code, material_description=description, material_group=group,
        # The unit is the material's, read from the workbook's material summary sheet;
        # a stock line only carries a copy of it.
        material=SimpleNamespace(material_code=code, description=description, uom="MT"),
        stock_qty=Decimal(qty), uom="MT", inventory_value_inr=Decimal(str(value_crore)) * Decimal("10000000"),
        stock_months=Decimal(str(months)),
    )


def _work_centre_stub():
    crore = Decimal("10000000")
    bands = {
        "critical_low_stock": [_centre_record(f"10000{index:04d}", "BARYTES API GRADE", "10", 0.4, 2.2, "310") for index in range(18)],
        "low_stock": [_centre_record("100002001", "BENTONITE OCMA GRADE", "10", 2.1, 2.65, "980")],
        "healthy_stock": [_centre_record("100002002", "CAUSTIC SODA FLAKES", "10", 4.5, 0.82, "140")],
        "slow_moving_stock": [_centre_record("100002003", "XANTHAN GUM POLYMER", "10", 9.2, 3.12, "60")],
        "excess_stock": [_centre_record("100002004", "DEFOAMER LIQUID", "10", 14.5, 1.89, "75")],
        "unclassified": [],
    }
    findings = [
        SimpleNamespace(exception_type="non_moving", details="Imported from Non Moving Inventory (row 5).",
                        material=SimpleNamespace(material_code="100002004", description="DEFOAMER LIQUID", uom="L"),
                        inventory_value_inr=Decimal("9700000"))
    ]
    return {
        "centre": SimpleNamespace(id=1, name="Ankleshwar DFS", zone="Western Onshore"),
        "thresholds": {"critical_low_stock_months": Decimal("1"), "low_stock_months": Decimal("3"), "slow_moving_months": Decimal("6"), "excess_stock_months": Decimal("12")},
        "groups": bands, "source_findings": findings, "selected_unit": None, "available_units": ["DFS", "ST"],
        "spec_groups": {
            key: [{"key": "DFC", "label": "Drilling Fluid Chemicals", "rows": rows, "total": len(rows), "omitted": 0}] if rows else []
            for key, rows in bands.items()
        },
        "source_spec_groups": [{"key": "DFC", "label": "Drilling Fluid Chemicals", "rows": findings, "total": len(findings), "omitted": 0}],
        "reporting_date": date(2026, 7, 31), "as_on_by_group": {"09": date(2026, 7, 31), "10": date(2026, 7, 31)},
        "previous_date": date(2026, 6, 30), "comparison_dates": [date(2026, 6, 30)],
        "comparison": {
            "previous_date": date(2026, 6, 30), "gap_days": 31, "common_materials": 4,
            "entrant_count": 7, "entrant_value": 15 * crore, "exit_count": 0, "exit_value": Decimal("0"),
            "like_for_like": 11 * crore, "like_for_like_prev": 9 * crore,
            "like_for_like_delta": 2 * crore, "like_for_like_change": 17.2,
        },
        "movers": {"up": [{"code": "090003008", "description": "OIL WELL CEMENT CLASS G HSR", "delta": Decimal("7560000")}], "down": []},
        "entrants": [], "exits": [],
        "kpis": {
            "total_value": 26 * crore, "prev_total_value": 9 * crore, "value_at_risk": 9 * crore, "at_risk_share": 37.0,
            "prev_value_at_risk": 6 * crore, "stockout_value": 10 * crore, "stockout_share": 39.5, "prev_stockout_value": 3 * crore,
            "line_count": 22, "material_count": 22, "portfolio_value": 870 * crore, "portfolio_share": 3.0,
            "source_case_count": 1, "source_case_value": Decimal("9700000"),
        },
        "health_mix": [{"key": key, "label": key.replace("_", " ").capitalize(), "value": 3 * crore, "count": len(rows), "share": 20.0, "prev": 2 * crore} for key, rows in bands.items() if rows],
        "top_materials": [{"code": "100002000", "description": "BARYTES API GRADE", "group": "10", "value": 4 * crore, "share": 15.8, "section": "Drilling Fluid Chemicals"}],
        "source_summary": [{"key": "non_moving", "label": "Non-moving materials", "count": 1, "value": Decimal("9700000")}],
    }


def test_work_centre_presentation_covers_the_page_registers(monkeypatch, tmp_path):
    from pptx import Presentation

    import app.core.services.inventory_monitoring as monitoring
    from app.core.services.inventory_presentation import build_work_centre_review_presentation

    monkeypatch.setattr(monitoring, "work_center_review_data", lambda work_center_id, unit=None, compare_date=None: _work_centre_stub())
    output, filename = build_work_centre_review_presentation(str(tmp_path), 1)

    assert filename == "Ankleshwar DFS Inventory Review 31 Jul 2026.pptx"
    titles = [
        shape.text_frame.text
        for slide in Presentation(output).slides
        for shape in slide.shapes
        if shape.has_text_frame and shape.text_frame.text
    ]
    assert "Ankleshwar DFS | Position summary" in titles
    assert "Critical low stock (1–13 of 18)" in titles
    assert "Low stock" in titles and "Slow-moving stock" in titles and "Excess stock" in titles
    assert "Non-moving, aged, surplus and transit cases" in titles
    assert "Decisions sought for Ankleshwar DFS" in titles
    assert not any("material group" in title.lower() for title in titles)
    assert any("4 materials were held at this centre on both" in title for title in titles)


def test_work_centre_presentation_requires_stock(monkeypatch, tmp_path):
    import pytest

    import app.core.services.inventory_monitoring as monitoring
    from app.core.services.inventory_presentation import build_work_centre_review_presentation

    empty = {**_work_centre_stub(), "reporting_date": None}
    empty["kpis"] = {**empty["kpis"], "line_count": 0}
    monkeypatch.setattr(monitoring, "work_center_review_data", lambda work_center_id, unit=None, compare_date=None: empty)
    with pytest.raises(ValueError):
        build_work_centre_review_presentation(str(tmp_path), 1)
