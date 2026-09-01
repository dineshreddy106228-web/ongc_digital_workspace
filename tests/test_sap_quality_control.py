"""Tests for the Corporate Chemistry SAP QC control-tower workflow."""
from __future__ import annotations

from datetime import date, datetime, timedelta
from io import BytesIO
import json
from pathlib import Path
import re
import sys
from types import SimpleNamespace

import pandas as pd
import pytest
from sqlalchemy import Integer

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from app import create_app
from app.extensions import db
from config import Config


def _xlsx(frame: pd.DataFrame, *, startrow: int = 0, title: str | None = None) -> bytes:
    buffer = BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        if title:
            pd.DataFrame([[title]]).to_excel(writer, index=False, header=False)
        frame.to_excel(writer, index=False, startrow=startrow)
    return buffer.getvalue()


def _inspection_export(*, plant: str = "10R2") -> bytes:
    return _xlsx(pd.DataFrame([
        ["890000005001", "00001234", plant, "01.08.2026", "", "REL CALC", ""],
        ["890000005002", "00005678", plant, "02.08.2026", "25.08.2026", "UD ICCO STUP", "A"],
        ["890000005003", "00009999", plant, "03.08.2026", "", "REL CALC", ""],
    ], columns=[
        "Inspection Lot", "Material", "Plant", "Start of Inspection", "End of Inspection",
        "System Status", "Usage Decision Code",
    ]))


def _notification_export(*, plant: str = "10R2", first_work_center: str = "MUDLAB") -> bytes:
    return _xlsx(pd.DataFrame([
        ["7001", "OPEN", "45000001", "10", "00001234", "Barytes", first_work_center, plant, "890000005001", "REL CALC", "01.08.2026", "20.08.2026", "", "6"],
        ["7002", "COMP", "45000002", "10", "00005678", "Calcium carbonate", "OILLAB", plant, "890000005002", "UD ICCO", "02.08.2026", "24.08.2026", "25.08.2026", "1"],
        ["7003", "OPEN", "45000003", "10", "00007777", "Unlinked item", "WATERLAB", plant, "0", "", "04.08.2026", "21.08.2026", "", "5"],
    ], columns=[
        "Notification No", "Notification Status", "Purchasing Document", "Item", "Material Number",
        "Material Description", "Work Center", "Plant", "Inspection Lot Number", "Status of Inspection Lot",
        "Start Date", "Planned End Date", "Completion Date", "Delay Days",
    ]), startrow=4, title="Date : 26.08.2026")


def _central_inspection_export() -> bytes:
    return _xlsx(pd.DataFrame([
        ["890000005101", "00001234", "10R2", "26.08.2026", "", "REL CALC", ""],
        ["890000005201", "00005678", "42R2", "26.08.2026", "", "REL CALC", ""],
    ], columns=[
        "Inspection Lot", "Material", "Plant", "Start of Inspection", "End of Inspection",
        "System Status", "Usage Decision Code",
    ]))


def _central_notification_export() -> bytes:
    return _xlsx(pd.DataFrame([
        ["7101", "OPEN", "45000001", "10", "00001234", "Barytes", "MUDLAB", "10R2", "890000005101", "REL CALC", "26.08.2026", "30.08.2026", "", "0"],
        ["7201", "OPEN", "45000002", "10", "00005678", "Xylene", "QUALILAB", "42R2", "890000005201", "REL CALC", "26.08.2026", "30.08.2026", "", "0"],
        ["7301", "OPEN", "45000003", "10", "00009999", "Glycol", "MUDLAB", "51R2", "890000005301", "REL CALC", "26.08.2026", "30.08.2026", "", "0"],
    ], columns=[
        "Notification No", "Notification Status", "Purchasing Document", "Item", "Material Number",
        "Material Description", "Work Center", "Plant", "Inspection Lot Number", "Status of Inspection Lot",
        "Start Date", "Planned End Date", "Completion Date", "Delay Days",
    ]), startrow=4, title="Date : 27.08.2026")


@pytest.fixture()
def sap_app(tmp_path):
    class _Config(Config):
        SQLALCHEMY_DATABASE_URI = f"sqlite:///{tmp_path / 'sap_qc.db'}"
        SQLALCHEMY_ENGINE_OPTIONS: dict = {}
        TESTING = True

    app = create_app(_Config)
    with app.app_context():
        from app.models.quality_control.qc_sample import QCSample
        from app.models.quality_control.qc_sap_monitoring import (
            QCNonSAPSample, QCNonSAPSampleUpdate, QCSAPLabUpdate, QCSAPMonitoringDisposition, QCSAPRecord,
            QCSAPSourceDocument, QCSAPUploadBatch,
        )

        # SQLite only auto-generates primary keys for an INTEGER column; the
        # production schema intentionally retains BIGINT primary keys for MySQL.
        for model in (
            QCSample, QCSAPSourceDocument, QCSAPUploadBatch, QCSAPRecord, QCSAPLabUpdate,
            QCSAPMonitoringDisposition, QCNonSAPSample, QCNonSAPSampleUpdate,
        ):
            model.__table__.c.id.type = Integer()
        db.create_all()
        yield app
        db.session.remove()


def test_native_sap_exports_are_parsed_and_unmatched_rows_are_retained():
    from app.core.services.sap_quality_control import (
        merge_sap_exports, parse_sap_inspection_workbook, parse_sap_notification_workbook,
    )

    inspections = parse_sap_inspection_workbook(_inspection_export(), "SAP_INSPECTION_20260826.xlsx")
    notifications = parse_sap_notification_workbook(_notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx")
    rows, reconciliation = merge_sap_exports(inspections.rows, notifications.rows)

    assert inspections.as_of_date == date(2026, 8, 26)
    assert notifications.as_of_date == date(2026, 8, 26)
    assert len(rows) == 4
    assert reconciliation == {"unmatched_inspection_count": 1, "unmatched_notification_count": 1}
    by_lot = {row["inspection_lot_number"]: row for row in rows if row["inspection_lot_number"]}
    assert by_lot["890000005001"]["official_status"] == "open"
    assert by_lot["890000005002"]["official_status"] == "completed"
    assert by_lot["890000005003"]["source_completeness"] == "inspection_lot_only"
    assert next(row for row in rows if row["notification_no"] == "7003")["source_completeness"] == "notification_only"


def test_sap_import_remains_authoritative_when_laboratory_logs_completion(sap_app):
    from flask import render_template
    from app.core.services.sap_quality_control import (
        create_sap_lab_update, import_sap_panvel_exports, sap_panvel_dashboard_data,
    )
    from app.models.quality_control.qc_sap_monitoring import QCSAPRecord

    batch = import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    db.session.commit()
    assert batch.record_count == 4
    record = QCSAPRecord.query.filter_by(inspection_lot_number="890000005001").one()
    assert record.official_status == "open"

    create_sap_lab_update(record.id, {
        "activity_status": "action_completed",
        "sampling_date": "2026-07-30",
        "actual_start_date": "2026-08-02",
        "expected_completion_date": "2026-08-27",
        "action_owner": "MUDLAB",
        "delay_reason": "Report sent for SAP posting",
        "update_note": "Awaiting QM update",
    }, None)
    db.session.commit()
    db.session.refresh(record)

    assert record.official_status == "open"
    data = sap_panvel_dashboard_data()
    assert data["kpis"]["open"] == 3
    assert data["kpis"]["accepted"] == 1
    assert data["kpis"]["rejected"] == 0
    assert data["kpis"]["awaiting_sap_confirmation"] == 1
    entry = next(item for item in data["records"] if item["record"].id == record.id)
    assert entry["reconciliation_key"] == "awaiting_sap_confirmation"
    assert entry["lab_update"].action_owner == "MUDLAB"
    assert entry["lab_update"].sampling_date == date(2026, 7, 30)
    assert entry["lab_update"].actual_start_date == date(2026, 8, 2)
    assert entry["courier_days"] == 2
    assert entry["time_in_queue_days"] == 1

    with sap_app.test_request_context("/quality-control/sap-control/labs/rgl_panvel"):
        page = render_template(
            "quality_control/sap_panvel_dashboard.html", can_control=False, **data,
        )
    assert "Sampling / lab start" in page
    assert "Sampling: 30 Jul 2026" in page
    assert "Notification: 01 Aug 2026" in page
    assert "Courier: 2 days" in page
    assert "Lab start: 02 Aug 2026" in page
    assert "Time in queue: 1 day" in page
    with sap_app.test_request_context("/quality-control/sap-control/labs/rgl_panvel"):
        laboratory_page = render_template(
            "quality_control/sap_panvel_dashboard.html",
            can_control=False, can_record_lab_updates=True, **data,
        )
    assert "Enter details" in laboratory_page
    assert "Update details" in laboratory_page
    assert "Date of sampling" in laboratory_page
    assert "Lab testing start date" in laboratory_page
    assert "Exclude non-actionable notification" not in laboratory_page


def test_lab_dashboard_hides_corporate_only_actions_from_a_reporting_laboratory(sap_app):
    """A reporting laboratory works its own list, not Corporate Chemistry's.

    Import, control-tower and all-labs actions belong to the corporate scope.
    A laboratory downloads its own deck.
    """
    from flask import render_template
    from app.core.services.sap_quality_control import import_sap_panvel_exports, sap_lab_dashboard_data

    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    db.session.commit()
    data = sap_lab_dashboard_data("rgl_panvel")

    with sap_app.test_request_context("/quality-control/sap-control/labs/rgl_panvel"):
        lab_page = render_template(
            "quality_control/sap_panvel_dashboard.html", can_control=False, **data,
        )
        corporate_page = render_template(
            "quality_control/sap_panvel_dashboard.html", can_control=True, **data,
        )

    for corporate_action in ("Import centre", "Upload at import centre", "All-labs presentation"):
        assert corporate_action not in lab_page
        assert corporate_action in corporate_page
    assert "Download presentation" in lab_page
    assert "Download presentation" in corporate_page
    # A laboratory may take its own SAP deck; it is not a corporate-only action.
    assert "Laboratory SAP deck" in lab_page
    assert "Laboratory SAP deck" in corporate_page
    # The operating-rule notice was removed; its snapshot provenance is not.
    assert "SAP is the final word" not in lab_page
    assert "SAP is the final word" not in corporate_page
    assert "plant 10R2" in lab_page
    assert "imported" in lab_page


def test_lab_navigator_dims_other_laboratories_and_reports_only_their_open_load(sap_app):
    """A laboratory sees the whole portfolio exists, and one number about it.

    Every location stays on the map so the reader knows the portfolio is
    larger than their bench, but a location they do not belong to is closed
    and carries only its actionable SAP-open count.
    """
    from app.core.services.quality_control import laboratory_landing_data, laboratory_navigator_data
    from app.core.services.sap_quality_control import (
        import_sap_panvel_exports, sap_open_counts_by_lab,
    )

    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    db.session.commit()

    open_counts = sap_open_counts_by_lab()
    assert open_counts["rgl_panvel"] == 3
    assert open_counts["rgl_vadodara"] == 0

    laboratories = laboratory_landing_data()
    # The navigator builds dashboard links, so it needs a request context.
    with sap_app.test_request_context("/quality-control/"):
        scoped = laboratory_navigator_data(
            laboratories, scope_lab_code="rgl_panvel", sap_open_counts=open_counts,
        )
        corporate = laboratory_navigator_data(
            laboratories, scope_lab_code=None, sap_open_counts=open_counts,
        )
    by_code = {entry["code"]: entry for entry in scoped}
    assert by_code["rgl_panvel"]["can_open"] is True
    assert by_code["rgl_panvel"]["sap_open_count"] == 3
    assert by_code["rgl_vadodara"]["can_open"] is False
    assert by_code["rgl_vadodara"]["sap_open_count"] == 0
    assert all(
        entry["can_open"] is (entry["code"] == "rgl_panvel") for entry in scoped
    ), "only the reader's own laboratory opens"

    assert all(entry["can_open"] for entry in corporate)


def test_usage_decision_reads_the_prefixed_sap_forms_without_over_matching():
    """SAP writes the decision as a bare code or with its UD prefix.

    The optional prefix and the word boundaries in this pattern were inert —
    written into a raw string as doubled backslashes, so they matched literal
    characters — which silently made every prefixed decision unrecorded.
    """
    from app.core.services.sap_quality_control import usage_decision_outcome

    assert usage_decision_outcome("A") == "accepted"
    assert usage_decision_outcome("R") == "rejected"
    assert usage_decision_outcome("UD A") == "accepted"
    assert usage_decision_outcome("UD R") == "rejected"
    assert usage_decision_outcome("UD  R") == "rejected"
    assert usage_decision_outcome("UDA") == "accepted"

    # A status string is not a decision, and a word merely starting with A or R
    # must not be read as one.
    for value in ("UD ICCO STUP", "ICCO", "REL CALC", "UD REJECTED", "ACCEPTED", "", None):
        assert usage_decision_outcome(value) is None, value


def test_derived_sap_readings_are_refreshed_on_every_write(sap_app):
    """The analytics group on these columns, so they must never go stale.

    They are materialised by a mapper event rather than at the importer, so a
    correction made anywhere re-derives them from the SAP values.
    """
    from app.core.services.sap_quality_control import import_sap_panvel_exports, latest_sap_batch
    from app.models.quality_control.qc_sap_monitoring import QCSAPRecord

    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    db.session.commit()

    # The import already carries a completed, accepted row with both dates.
    completed = QCSAPRecord.query.filter_by(
        lab_code="rgl_panvel", notification_no="7002",
    ).one()
    assert completed.usage_outcome == "accepted"
    assert completed.turnaround_days == (date(2026, 8, 25) - date(2026, 8, 2)).days

    # A correction outside the importer re-derives rather than leaving a
    # reading that disagrees with the code and dates beside it.
    completed.usage_decision_code = "UD R"
    completed.completion_date = date(2026, 8, 10)
    db.session.commit()
    assert completed.usage_outcome == "rejected"
    assert completed.turnaround_days == 8

    # A completion recorded before the start is a source fault, not zero days.
    completed.completion_date = date(2026, 7, 1)
    db.session.commit()
    assert completed.turnaround_days is None


def test_portfolio_analytics_counts_in_one_grouped_query(sap_app):
    """The database does the counting; Python only maps the register."""
    from sqlalchemy import event
    from app.core.services.sap_quality_control import (
        import_sap_panvel_exports, import_sap_lab_exports, sap_portfolio_analytics,
    )

    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    import_sap_lab_exports(
        "rgl_vadodara", _inspection_export(plant="23R2"), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(plant="23R2"), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    db.session.commit()

    statements = []

    def record_statement(_conn, _cursor, statement, *_args):
        if statement.lstrip().upper().startswith("SELECT"):
            statements.append(statement)

    event.listen(db.engine, "after_cursor_execute", record_statement)
    try:
        analytics = sap_portfolio_analytics()
    finally:
        event.remove(db.engine, "after_cursor_execute", record_statement)

    record_reads = [
        statement for statement in statements if "qc_sap_records" in statement.lower()
    ]
    assert len(record_reads) == 1, record_reads
    assert "group by" in record_reads[0].lower()
    # It never grows with the number of laboratories in scope.
    assert analytics["totals"]["total"] == 8


def test_portfolio_analytics_reads_the_whole_recorded_load_not_one_snapshot(sap_app, monkeypatch):
    """The analysis counts every record ever imported, and rates honestly.

    A completed sample drops out of the daily snapshot, so a snapshot-only
    analysis would never see an outcome. Rates are taken over decided samples,
    and a material needs enough decisions before it is ranked on a percentage.
    """
    from app.core.services import corporate_specifications
    from app.core.services.sap_quality_control import (
        MIN_DECISIONS_FOR_RATE, import_sap_panvel_exports, import_sap_lab_exports,
        latest_sap_batch, sap_portfolio_analytics,
    )
    from app.models.quality_control.qc_sap_monitoring import QCSAPRecord

    monkeypatch.setattr(corporate_specifications, "catalogue", lambda: [
        {
            "material_code": "1234", "category": "DFC", "category_label": "Drilling Fluid Chemicals",
            "spec_number": "ONGC / DFC / 01 / 2026", "chemical_name": "Barytes",
            "standard_days": "10", "on_register": True,
        },
    ])
    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    import_sap_lab_exports(
        "rgl_vadodara", _inspection_export(plant="23R2"), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(plant="23R2"), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    db.session.commit()

    # Barytes is decided often enough to be ranked; one of those is a rejection.
    barytes = QCSAPRecord.query.filter_by(
        lab_code="rgl_panvel", material_description="Barytes",
    ).one()
    barytes.official_status = "completed"
    barytes.usage_decision_code = "R"
    barytes.completion_date = date(2026, 8, 20)
    for index in range(MIN_DECISIONS_FOR_RATE - 1):
        db.session.add(QCSAPRecord(
            source_key=f"extra-barytes-{index}", lab_code="rgl_panvel",
            material_code="00001234", material_description="Barytes",
            official_status="completed", usage_decision_code="A",
            start_inspection_date=date(2026, 8, 1), completion_date=date(2026, 8, 9),
            last_seen_batch_id=latest_sap_batch("rgl_panvel").id,
        ))
    db.session.commit()

    analytics = sap_portfolio_analytics()

    # Eight rows exist across both laboratories plus the four added completions.
    assert analytics["has_data"] is True
    assert analytics["totals"]["total"] == 12
    assert analytics["totals"]["rejected"] == 1
    assert analytics["totals"]["accepted"] == 6
    # Rates use decided samples only, never the open ones.
    assert analytics["totals"]["decided"] == 7
    assert analytics["totals"]["rejection_rate"] == round(1 / 7 * 100, 1)

    by_load = {item["material_description"]: item for item in analytics["materials_by_load"]}
    assert analytics["materials_by_load"][0]["material_description"] == "Barytes"
    assert by_load["Barytes"]["total"] == 6
    assert by_load["Barytes"]["subgroup_label"] == "Drilling Fluid Chemicals"

    # Only Barytes clears the minimum-decisions bar, so nothing else is ranked
    # on a percentage taken from one or two samples.
    ranked = analytics["materials_by_failure"]
    assert [item["material_description"] for item in ranked] == ["Barytes"]
    assert ranked[0]["decided"] == MIN_DECISIONS_FOR_RATE
    assert ranked[0]["rejection_rate"] == 20.0

    # Turnaround is measured against SAP dates. Four added samples ran 8 days
    # against a 10-day standard; the rejected one ran 19 and missed it, so the
    # laboratory reads 4 of 5 within STT rather than a flat pass.
    panvel = next(
        item for item in analytics["laboratories"]
        if item["laboratory"]["code"] == "rgl_panvel"
    )
    assert panvel["median_turnaround_days"] == 8.0
    assert (panvel["within_stt"], panvel["stt_measured"]) == (4, 5)
    assert panvel["stt_on_time_rate"] == 80.0
    # A laboratory that has completed nothing measurable sorts last rather
    # than ranking as a perfect score.
    assert analytics["laboratories"][0]["laboratory"]["code"] == "rgl_panvel"
    assert analytics["laboratories"][-1]["stt_on_time_rate"] is None


def test_non_sap_register_summarises_without_touching_the_sap_position(sap_app):
    """Management can now see the non-SAP load; it still is not SAP.

    The register is reported on its own terms — a declared result, not a usage
    decision — and none of it reaches the SAP figures beside it.
    """
    from app.core.services.sap_quality_control import (
        import_sap_panvel_exports, non_sap_register_data, sap_management_data,
    )
    from app.models.quality_control.qc_sap_monitoring import QCNonSAPSample

    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    today = date.today()
    db.session.add_all([
        QCNonSAPSample(
            lab_code="rgl_panvel", sample_reference="NS-1", chemical_name="Produced water",
            current_status="under_testing", sample_receipt_date=today - timedelta(days=10),
            expected_completion_date=today - timedelta(days=2),
        ),
        QCNonSAPSample(
            lab_code="rgl_panvel", sample_reference="NS-2", chemical_name="Produced water",
            current_status="awaiting_sample", sample_receipt_date=today - timedelta(days=3),
        ),
        QCNonSAPSample(
            lab_code="rgl_vadodara", sample_reference="NS-3", chemical_name="Glycol",
            current_status="closed_fail", reported_outcome="fail",
            expected_completion_date=today - timedelta(days=30),
        ),
        QCNonSAPSample(
            lab_code="rgl_vadodara", sample_reference="NS-4", chemical_name="Glycol",
            current_status="closed_pass", reported_outcome="pass",
        ),
    ])
    db.session.commit()

    register = non_sap_register_data()
    kpis = register["non_sap_kpis"]

    assert kpis["total"] == 4
    assert kpis["pending"] == 2
    # A closed sample is history, so a past expected date does not make it late.
    assert kpis["overdue"] == 1
    assert kpis["no_expected_date"] == 1
    # The rate is over declared closures only, never over pending work.
    assert (kpis["closed_pass"], kpis["closed_fail"], kpis["closed"]) == (1, 1, 2)
    assert kpis["fail_rate"] == 50.0

    # Overdue work leads and closed work sorts to the end; within each band the
    # nearer committed date comes first, and no committed date comes last.
    assert [item["sample"].sample_reference for item in register["non_sap_entries"]] == [
        "NS-1", "NS-2", "NS-3", "NS-4",
    ]
    assert [item["laboratory"]["name"] for item in register["non_sap_by_laboratory"]] == [
        "RGL Panvel", "RGL Vadodara",
    ]
    chemicals = {item["chemical_name"]: item for item in register["non_sap_chemicals"]}
    assert chemicals["Produced water"]["pending"] == 2
    assert chemicals["Glycol"]["failed"] == 1

    # None of it has moved into the SAP position.
    management = sap_management_data()
    assert management["kpis"]["total"] == 4
    assert management["kpis"]["rejected"] == 0
    assert all(
        entry["record"].material_description != "Produced water"
        for entry in management["action_entries"]
    )


def test_import_centre_absorbed_the_control_tower(sap_app):
    """One SAP operations screen, carrying everything the tower uniquely held.

    The control tower duplicated this page's central import form. What it did
    not duplicate was the per-laboratory position table, the cross-laboratory
    non-SAP register, and the link to each laboratory's own SAP deck — so
    removing the tower had to bring those three here rather than drop them.
    """
    from flask import render_template
    from app.core.services.quality_control import (
        current_monitoring_day, laboratory_landing_data, laboratory_navigator_data,
    )
    from app.core.services.sap_quality_control import sap_control_data

    with sap_app.test_request_context("/quality-control/"):
        laboratories = laboratory_landing_data()
        monitoring_day = current_monitoring_day()
        landing = render_template(
            "quality_control/landing.html",
            laboratories=laboratories,
            designated_laboratories=[],
            mapped_laboratory_total=len(laboratories),
            map_laboratories=laboratory_navigator_data(laboratories, monitoring_day["date"]),
            monitoring_day=monitoring_day, is_superuser=True,
        )
        sap_data = sap_control_data()
        import_centre = render_template(
            "quality_control/data_import.html",
            laboratories=[], monitoring_day=monitoring_day,
            sap_control_cards=sap_data["control_cards"],
            sap_laboratories=sap_data["sap_laboratories"],
            sap_plant_mappings=sap_data["sap_plant_mappings"],
            all_laboratories=sap_data["all_laboratories"],
            non_sap_entries=sap_data["non_sap_entries"],
            non_sap_statuses=sap_data["non_sap_statuses"],
            non_sap_status_labels=sap_data["non_sap_status_labels"],
            can_control=True,
        )

    # No tile offers a page that no longer exists.
    tiles = re.findall(r'<a class="mod-tile" href="([^"]+)".*?<strong>([^<]+)</strong>', landing)
    tiles_by_name = {name: href for href, name in tiles}
    assert "SAP Control Tower" not in tiles_by_name
    assert tiles_by_name["Import Centre"] == "/quality-control/data-import"
    assert "/quality-control/sap-control" not in landing

    # The import screen keeps its own identity and the tower's three sections.
    assert "SAP Import Centre" in import_centre
    heading = re.search(r'<h1 class="mod-page-title">.*?</h1>', import_centre, re.S).group()
    assert "Control Tower" not in heading
    assert "RGL and IDWE SAP position" in import_centre
    assert "Exclusion review" in import_centre
    assert 'id="non-sap-register"' in import_centre
    assert "Add a non-SAP sample returned by a laboratory" in import_centre
    # Each row opens the laboratory's own control view, where records are
    # excluded, reinstated and updated. The deck button beside it needs a
    # snapshot, so it is asserted where one exists, on the lab dashboard.
    assert "Open control view" in import_centre


def test_panvel_import_rejects_a_non_panvel_plant():
    from app.core.services.sap_quality_control import parse_sap_notification_workbook

    with pytest.raises(ValueError, match="plant 10R2 only"):
        parse_sap_notification_workbook(_notification_export(plant="11A1"), "SAP_NOTIFICATIONS_20260826.xlsx")


def test_central_sap_import_splits_rows_by_the_approved_plant_map(sap_app):
    from app.core.services.sap_quality_control import (
        SAP_PLANT_LAB_CODES,
        import_central_sap_exports,
    )
    from app.models.quality_control.qc_sap_monitoring import QCSAPRecord

    batches = import_central_sap_exports(
        _central_inspection_export(), "SAP_INSP_20260827.xlsx",
        _central_notification_export(), "SAP_ZLABIMS_20260827.xlsx", None,
    )
    db.session.commit()

    assert SAP_PLANT_LAB_CODES == {
        "10R2": "rgl_panvel", "23R2": "rgl_vadodara", "42R2": "rgl_chennai",
        "41B1": "rgl_rajahmundry", "50R2": "rgl_jorhat_sivasagar", "51R2": "rgl_jorhat",
        "70T1": "idwe_dehradun",
    }
    assert {(batch.plant_code, batch.lab_code) for batch in batches} == {
        ("10R2", "rgl_panvel"), ("42R2", "rgl_chennai"), ("51R2", "rgl_jorhat"),
    }
    jorhat_record = QCSAPRecord.query.filter_by(lab_code="rgl_jorhat").one()
    assert jorhat_record.source_completeness == "notification_only"
    assert jorhat_record.plant_code == "51R2"


def test_central_sap_import_rejects_a_wholly_foreign_export_before_persisting(sap_app):
    """Rows from plants no laboratory owns are set aside; a report made only of
    them is a different thing, and still stops the upload with nothing written."""
    from app.core.services.sap_quality_control import import_central_sap_exports
    from app.models.quality_control.qc_sap_monitoring import QCSAPRecord, QCSAPUploadBatch

    with pytest.raises(ValueError, match="no rows for an RGL or IDWE plant"):
        import_central_sap_exports(
            _inspection_export(plant="50R1"), "SAP_INSP_20260827.xlsx",
            _notification_export(plant="50R1"), "SAP_ZLABIMS_20260827.xlsx", None,
        )
    assert QCSAPUploadBatch.query.count() == 0
    assert QCSAPRecord.query.count() == 0


def test_import_explains_when_a_notification_export_is_selected_as_inspection_lots():
    from app.core.services.sap_quality_control import parse_sap_inspection_workbook

    with pytest.raises(ValueError, match="Notifications / ZLABIMS export"):
        parse_sap_inspection_workbook(_notification_export(plant="42R2"), "RGL_Chennai_Inspection_Lots.xlsx", expected_plant=None)


def test_other_sap_reporting_lab_accepts_its_single_plant_snapshot(sap_app):
    from app.core.services.sap_quality_control import import_sap_lab_exports, sap_lab_dashboard_data

    batch = import_sap_lab_exports(
        "rgl_vadodara", _inspection_export(plant="23R2"), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(plant="23R2"), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    db.session.commit()

    assert batch.lab_code == "rgl_vadodara"
    assert batch.plant_code == "23R2"
    data = sap_lab_dashboard_data("rgl_vadodara")
    assert data["laboratory"]["name"] == "RGL Vadodara"
    assert data["kpis"]["open"] == 3


def test_module_sample_register_uses_current_sap_records_not_weekly_workbook_rows(sap_app, monkeypatch):
    from flask import render_template
    from app.core.services import corporate_specifications
    from app.core.services.sap_quality_control import (
        import_sap_panvel_exports, sap_panvel_dashboard_data, sap_sample_register_data,
    )
    from app.models.quality_control.qc_sample import QCSample

    monkeypatch.setattr(corporate_specifications, "catalogue", lambda: [
        {
            "material_code": "1234", "category": "DFC", "category_label": "Drilling Fluid Chemicals",
            "spec_number": "ONGC / DFC / 01 / 2026", "chemical_name": "Barytes", "on_register": True,
        },
        {
            "material_code": "00005678", "category": "PC", "category_label": "Production Chemicals",
            "spec_number": "ONGC / PC / 02 / 2026", "chemical_name": "Calcium carbonate", "on_register": True,
        },
    ])
    batch = import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    # This is deliberately a plausible old register row. The shared register
    # must not read it after the SAP source migration.
    db.session.add(QCSample(
        source_key="legacy-weekly-row", lab_code="rgl_panvel",
        chemical_name="Legacy weekly workbook sample", result_status="under_testing",
    ))
    db.session.commit()

    register = sap_sample_register_data()
    entries = register["entries"]

    assert len(entries) == 4
    assert {entry["batch"].id for entry in entries} == {batch.id}
    assert all(entry["record"].material_description != "Legacy weekly workbook sample" for entry in entries)
    assert len(sap_sample_register_data(status="open")["entries"]) == 3
    assert len(sap_sample_register_data(status="completed")["entries"]) == 1
    assert len(sap_sample_register_data(status="accepted")["entries"]) == 1
    assert [entry["record"].notification_no for entry in sap_sample_register_data(search="MUDLAB")["entries"]] == ["7001"]
    barytes = next(entry for entry in entries if entry["record"].material_description == "Barytes")
    assert barytes["subgroup_label"] == "Drilling Fluid Chemicals"
    assert barytes["specification_no"] == "ONGC / DFC / 01 / 2026"
    assert barytes["specification_match"] is True
    assert [entry["record"].material_description for entry in sap_sample_register_data(subgroup="PC")["entries"]] == [
        "Calcium carbonate",
    ]
    assert len(sap_sample_register_data(subgroup="not_in_corporate_specification")["entries"]) == 2
    # The register is read category by category, so rows arrive grouped in
    # Corporate Specifications order with unmatched material codes collected
    # in a final section rather than spread through the table.
    groups = register["groups"]
    assert [(group["key"], len(group["entries"])) for group in groups] == [
        ("DFC", 1), ("PC", 1), ("not_in_corporate_specification", 2),
    ]
    assert [group["label"] for group in groups] == [
        "Drilling Fluid Chemicals", "Production Chemicals", "Not in Corporate Specification",
    ]
    assert [group["is_unmatched"] for group in groups] == [False, False, True]
    assert sum(len(group["entries"]) for group in groups) == len(entries)
    dashboard = sap_panvel_dashboard_data()
    matched_record = next(entry["record"] for entry in entries if entry["record"].notification_no == "7001")
    notification_only_record = next(entry["record"] for entry in entries if entry["record"].notification_no == "7003")
    assert matched_record.start_inspection_date == date(2026, 8, 1)
    assert notification_only_record.start_inspection_date is None
    assert notification_only_record.notification_start_date == date(2026, 8, 4)
    assert next(
        entry for entry in dashboard["records"] if entry["record"].material_description == "Barytes"
    )["subgroup_label"] == "Drilling Fluid Chemicals"
    with sap_app.test_request_context("/quality-control/history"):
        page = render_template(
            "quality_control/samples.html",
            filters={"lab": "", "search": "", "status": "", "subgroup": ""},
            **register,
        )
    assert "SAP Sample Register" in page
    assert "Drilling Fluid Chemicals" in page
    assert "Production Chemicals" in page
    assert "Not in Corporate Specification" in page
    assert "Legacy weekly workbook sample" not in page
    assert "SAP receipt 01 Aug 2026" in page
    assert page.index("Drilling Fluid Chemicals") < page.index("Production Chemicals") < page.index("Not in Corporate Specification")
    assert "qc-sap-subgroup-heading" in page
    assert "Notification date 04 Aug 2026" in page
    with sap_app.test_request_context("/quality-control/sap-control/labs/rgl_panvel"):
        dashboard_page = render_template(
            "quality_control/sap_panvel_dashboard.html", can_control=False, **dashboard,
        )
    assert "Corporate Specification sub-group" in dashboard_page
    assert "Material / specification" in dashboard_page
    assert dashboard_page.count("Drilling Fluid Chemicals") == 1
    assert "SAP receipt 01 Aug 2026" in dashboard_page
    assert "Notification date 04 Aug 2026" in dashboard_page
    assert "Complete every follow-up item" in dashboard_page
    assert "contact <strong>Corporate Chemistry</strong>" in dashboard_page


def test_sap_monitoring_uses_corporate_stt_not_sap_planned_end(sap_app, monkeypatch):
    from app.core.services import corporate_specifications
    from app.core.services.sap_quality_control import (
        import_sap_panvel_exports, sap_lab_dashboard_data, sap_management_data,
    )

    monkeypatch.setattr(corporate_specifications, "catalogue", lambda: [{
        "material_code": "1234", "category": "DFC", "category_label": "Drilling Fluid Chemicals",
        "spec_number": "ONGC / DFC / 01 / 2026", "chemical_name": "Barytes",
        "standard_days": 30, "on_register": True,
    }])
    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    db.session.commit()

    dashboard = sap_lab_dashboard_data("rgl_panvel")
    barytes = next(item for item in dashboard["records"] if item["record"].notification_no == "7001")

    # SAP's planned end for this row is 20 Aug.  It must not make the row
    # overdue: its Corporate Specification STT runs for 30 days from receipt.
    assert barytes["record"].planned_end_date == date(2026, 8, 20)
    assert barytes["stt_days"] == 30
    assert barytes["stt_due_date"] == date(2026, 8, 31)
    assert barytes["stt_overdue"] is False
    assert dashboard["kpis"]["stt_overdue"] == 0
    assert "planned_overdue" not in dashboard["kpis"]

    management = sap_management_data({"rgl_panvel"})
    assert management["kpis"]["stt_overdue"] == 0
    assert "planned_overdue" not in management["kpis"]


def test_non_sap_register_is_separate_from_the_sap_position(sap_app):
    from flask import render_template
    from app.core.services.sap_quality_control import (
        create_non_sap_sample, import_sap_panvel_exports, sap_panvel_dashboard_data,
        update_non_sap_sample,
    )

    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    sample = create_non_sap_sample({
        "lab_code": "rgl_panvel", "sample_reference": "PANVEL-LOCAL-08", "chemical_name": "Field control sample",
        "current_status": "under_testing", "expected_completion_date": "2026-08-28", "action_owner": "QC Panvel",
        "delay_reason": "Awaiting calibration", "update_note": "Declared separately from SAP",
    }, None)
    db.session.commit()
    assert sap_panvel_dashboard_data()["kpis"]["open"] == 3
    assert sap_panvel_dashboard_data()["kpis"]["non_sap_pending"] == 1
    non_sap_allocation = next(
        item for item in sap_panvel_dashboard_data()["work_centers"]
        if item["is_non_sap"]
    )
    assert non_sap_allocation == {
        "name": "Non-SAP samples", "open": 1, "stt_overdue": 0,
        "awaiting_lab": 0, "is_non_sap": True,
    }
    with sap_app.test_request_context("/quality-control/sap-control/labs/rgl_panvel"):
        page = render_template(
            "quality_control/sap_panvel_dashboard.html",
            can_control=False, can_record_lab_updates=True,
            **sap_panvel_dashboard_data(),
        )
    assert "Add non-SAP sample" in page
    assert "Field control sample" in page
    assert "Save non-SAP update" in page
    assert "Non-SAP samples" in page
    assert "Laboratory-declared · not an SAP work center" in page
    assert "/sap-control/labs/rgl_panvel/non-sap" in page

    with pytest.raises(ValueError, match="does not belong to this laboratory"):
        update_non_sap_sample(sample.id, {"current_status": "under_testing"}, None, lab_code="rgl_vadodara")

    update_non_sap_sample(sample.id, {
        "current_status": "closed_pass", "expected_completion_date": "2026-08-28", "action_owner": "QC Panvel",
        "delay_reason": "", "update_note": "Reported accepted", "reported_outcome": "pass",
    }, None)
    db.session.commit()
    data = sap_panvel_dashboard_data()
    assert data["kpis"]["open"] == 3
    assert data["kpis"]["non_sap_pending"] == 0


def test_qc_admin_exclusion_is_auditable_and_reopens_for_sap_assignment_change(sap_app):
    from app.core.services.sap_quality_control import (
        exclude_sap_record_from_monitoring, import_sap_panvel_exports,
        reinstate_sap_record_for_monitoring, sap_panvel_dashboard_data,
    )
    from app.models.quality_control.qc_sap_monitoring import QCSAPRecord

    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    db.session.commit()
    record = QCSAPRecord.query.filter_by(notification_no="7001").one()
    exclusion = exclude_sap_record_from_monitoring(record.id, {
        "exclusion_reason": "junk_notification",
        "exclusion_note": "Created in error; no laboratory work was requested.",
    }, None)
    db.session.commit()

    data = sap_panvel_dashboard_data()
    assert exclusion.decision == "exclude_non_actionable"
    assert data["kpis"]["open"] == 2
    assert data["kpis"]["excluded_from_monitoring"] == 1
    assert all(item["record"].id != record.id for item in data["open_records"])
    assert data["excluded_entries"][0]["exclusion_reason_label"] == "Junk / test notification"

    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(first_work_center="NEWLAB"), "SAP_NOTIFICATIONS_20260827.xlsx", None,
    )
    db.session.commit()
    changed = sap_panvel_dashboard_data()
    assert changed["kpis"]["excluded_from_monitoring"] == 0
    assert changed["kpis"]["exclusion_review"] == 1
    assert changed["exclusion_review_entries"][0]["record"].id == record.id

    reinstate_sap_record_for_monitoring(record.id, None)
    db.session.commit()
    reinstated = sap_panvel_dashboard_data()
    assert reinstated["kpis"]["open"] == 3
    assert reinstated["kpis"]["exclusion_review"] == 0


def test_management_data_uses_latest_sap_records_and_excludes_non_sap_rows(sap_app):
    from app.core.services.sap_quality_control import (
        create_non_sap_sample, exclude_sap_record_from_monitoring,
        import_sap_panvel_exports, sap_management_data,
    )
    from app.models.quality_control.qc_sap_monitoring import QCSAPRecord

    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    create_non_sap_sample({
        "lab_code": "rgl_panvel", "sample_reference": "LOCAL-01",
        "chemical_name": "Not in SAP", "current_status": "under_testing",
    }, None)
    db.session.commit()

    data = sap_management_data()
    assert data["reporting_labs"] == 1
    assert data["kpis"]["total"] == 4
    assert data["kpis"]["actionable_open"] == 3
    assert data["kpis"]["accepted"] == 1
    assert len(data["action_entries"]) == 3
    assert all(item["record"].material_description != "Not in SAP" for item in data["action_entries"])

    record = QCSAPRecord.query.filter_by(notification_no="7001").one()
    exclude_sap_record_from_monitoring(record.id, {
        "exclusion_reason": "junk_notification", "exclusion_note": "Test entry",
    }, None)
    db.session.commit()
    excluded = sap_management_data()
    assert excluded["kpis"]["total"] == 4
    assert excluded["kpis"]["actionable_open"] == 2
    assert excluded["kpis"]["excluded"] == 1


def test_management_separates_completed_notification_only_records_without_ud_details(sap_app):
    from flask import render_template
    from app.core.services.qc_management_charts import management_chart_series
    from app.core.services.sap_quality_control import (
        import_sap_panvel_exports, non_sap_register_data, sap_management_data,
        sap_portfolio_analytics,
    )
    from app.models.quality_control.qc_sap_monitoring import QCSAPRecord

    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    notification_only = QCSAPRecord.query.filter_by(notification_no="7003").one()
    notification_only.inspection_lot_number = "890000004999"
    notification_only.completion_date = date(2026, 8, 20)
    notification_only.official_status = "completed"
    db.session.commit()

    data = sap_management_data()

    assert data["kpis"]["completed_without_ud_details"] == 1
    assert [item["record"].notification_no for item in data["completed_without_ud_entries"]] == ["7003"]
    assert data["completed_without_ud_entries"][0]["record"].inspection_lot_number == "890000004999"
    non_sap = non_sap_register_data()
    portfolio = sap_portfolio_analytics()
    with sap_app.test_request_context("/quality-control/management-review"):
        page = render_template(
            "quality_control/portfolio_management_review.html",
            portfolio=portfolio,
            charts=management_chart_series(data, portfolio, non_sap),
            **data, **non_sap,
        )
    assert "Completed without inspection-lot / UD details" in page
    assert "No paired QA33 lot row · no UD details" in page
    assert "890000004999" in page


def test_sap_management_presentation_uses_the_current_snapshot(sap_app):
    from app.core.services.qc_presentation import (
        build_sap_panvel_presentation, build_sap_portfolio_management_presentation,
    )
    from app.core.services.sap_quality_control import import_sap_panvel_exports

    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    db.session.commit()
    output, filename = build_sap_panvel_presentation(sap_app.static_folder)

    assert filename == "RGL Panvel SAP QC Presentation 26 Aug 2026.pptx"
    assert output.read(2) == b"PK"

    portfolio_output, portfolio_filename = build_sap_portfolio_management_presentation(sap_app.static_folder)
    assert portfolio_filename == "QC SAP Management Review 26 Aug 2026.pptx"
    assert portfolio_output.read(2) == b"PK"

    single_lab_output, single_lab_filename = build_sap_portfolio_management_presentation(
        sap_app.static_folder, {"rgl_panvel"},
    )
    assert single_lab_filename == "RGL Panvel Management Review 26 Aug 2026.pptx"
    assert single_lab_output.read(2) == b"PK"

def test_lab_deck_lists_unassigned_notifications_for_an_active_inactive_call(sap_app):
    """A notification SAP never routed to a work center gets its own page.

    Without a work center nobody at the laboratory recognises the item as
    theirs, so it can sit open indefinitely while still counting against the
    bench. The deck already lists it among the action items; this page asks
    the one question that clears it — is it still live?
    """
    from pptx import Presentation
    from app.core.services.qc_presentation import build_sap_lab_presentation
    from app.core.services.sap_quality_control import (
        import_sap_panvel_exports, sap_lab_dashboard_data,
    )

    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(first_work_center=""), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    db.session.commit()

    data = sap_lab_dashboard_data("rgl_panvel")
    unassigned = [
        item for item in data["open_records"]
        if not item["record"].work_center and item["record"].notification_no
    ]
    assigned = [item for item in data["open_records"] if item["record"].work_center]
    # An inspection lot carrying no notification is not something a laboratory
    # can call active or inactive, so it stays off this page.
    lot_only = [
        item for item in data["open_records"]
        if not item["record"].work_center and not item["record"].notification_no
    ]
    assert unassigned, "the fixture must leave at least one open notification unrouted"
    assert assigned, "and at least one routed, so the page is not simply everything"
    assert lot_only, "and at least one lot-only record, which must be excluded"

    output, _ = build_sap_lab_presentation("rgl_panvel", sap_app.static_folder)
    pages = []
    for slide in Presentation(output).slides:
        text = "\n".join(
            shape.text_frame.text for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        )
        if "Unassigned notifications" in text:
            tables = [shape.table for shape in slide.shapes if shape.has_table]
            pages.append((text, tables[0] if tables else None))

    assert len(pages) == 1
    text, table = pages[0]
    assert table is not None
    assert f"(1\u2013{len(unassigned)} of {len(unassigned)})" in text

    header = [cell.text.strip() for cell in table.rows[0].cells]
    assert "Active / Inactive" in header
    assert "Laboratory remark" in header

    # Every unrouted record, and only those.
    body = [" | ".join(cell.text for cell in row.cells) for row in list(table.rows)[1:]]
    assert len(body) == len(unassigned)
    for item in unassigned:
        reference = item["record"].notification_no or item["record"].inspection_lot_number
        assert any(reference in line for line in body)
    for item in assigned:
        assert not any(item["record"].work_center in line for line in body)
    for item in lot_only:
        assert not any(item["record"].inspection_lot_number in line for line in body)


def test_lab_deck_omits_the_unassigned_page_when_sap_routed_everything(sap_app):
    """No unrouted record, no page — the deck does not carry an empty ask."""
    from pptx import Presentation
    from app.core.services.qc_presentation import build_sap_lab_presentation
    from app.core.services.sap_quality_control import (
        import_sap_panvel_exports, sap_lab_dashboard_data,
    )

    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    db.session.commit()
    for item in sap_lab_dashboard_data("rgl_panvel")["open_records"]:
        if not item["record"].work_center:
            item["record"].work_center = "MUDLAB"
    db.session.commit()
    assert not [
        item for item in sap_lab_dashboard_data("rgl_panvel")["open_records"]
        if not item["record"].work_center and item["record"].notification_no
    ]

    output, _ = build_sap_lab_presentation("rgl_panvel", sap_app.static_folder)
    for slide in Presentation(output).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                assert "Unassigned notifications" not in shape.text_frame.text


def test_register_reports_matches_before_the_display_cap(sap_app, monkeypatch):
    """The caption counts what matched, not what fitted on the page.

    It used to be measured on the truncated list, so a filter matching
    thousands and a filter matching exactly the cap both read the same. There
    was no way to tell the register was holding anything back.
    """
    from app.core.services import sap_quality_control
    from app.core.services.sap_quality_control import (
        import_sap_panvel_exports, sap_sample_register_data,
    )

    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    db.session.commit()

    uncapped = sap_sample_register_data()
    assert uncapped["total_matching"] == len(uncapped["entries"])
    assert uncapped["is_truncated"] is False

    # Below the cap the two agree, so force the cap under the fixture instead
    # of importing hundreds of rows to reach it.
    monkeypatch.setattr(sap_quality_control, "SAP_REGISTER_VISIBLE_LIMIT", 2)
    capped = sap_sample_register_data()
    assert capped["total_matching"] == uncapped["total_matching"]
    assert len(capped["entries"]) == 2
    assert capped["is_truncated"] is True
    assert capped["visible_limit"] == 2
    # The grouped view must not out-count the rows it was built from.
    assert sum(len(group["entries"]) for group in capped["groups"]) == 2


def test_pairing_label_says_what_is_missing_not_which_export_it_came_from():
    """"Notification Only" sat directly under the lot number it seemed to deny.

    The label described provenance, so a row showing a bold inspection lot was
    captioned as notification-only and read as self-contradictory. It has to
    say the thing the reader needs instead: that stated lot is absent from the
    paired QA33 export, which is why no usage decision could be joined.
    """
    from app.core.services.sap_quality_control import source_completeness_label
    from app.models.quality_control.qc_sap_monitoring import QCSAPRecord

    def record(completeness, lot=None):
        return QCSAPRecord(source_completeness=completeness, inspection_lot_number=lot)

    stated = source_completeness_label(record("notification_only", "890000040483"))
    assert stated == "Notification only · stated lot not in QA33 export"

    # A notification naming no lot has nothing missing from QA33 to chase.
    assert source_completeness_label(record("notification_only")) == (
        "Notification only · no inspection lot stated"
    )
    assert source_completeness_label(record("inspection_lot_only", "890000040432")) == (
        "Inspection lot only · no notification in SAP"
    )
    assert source_completeness_label(record("matched", "890000040435")) == (
        "Matched · lot and notification paired"
    )
    # An unrecognised value still renders as something readable.
    assert source_completeness_label(record("some_new_state")) == "Some New State"


def test_register_export_splits_by_subgroup_and_exports_every_match(sap_app):
    """The download is the uncapped answer the page cannot render.

    Work-centre records go to their Corporate Specification sub-group sheet,
    because that is the axis a chemist reads along. The two source sheets are
    the exceptions worth chasing, laboratory-wise and newest first.
    """
    from openpyxl import load_workbook
    from app.core.services.qc_register_export import build_sample_register_workbook
    from app.core.services.sap_quality_control import (
        import_sap_panvel_exports, sap_sample_register_data,
    )

    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    db.session.commit()

    entries = sap_sample_register_data(limit=None)["entries"]
    with_centre = [item for item in entries if item["record"].work_center]
    lot_only = [item for item in entries
                if item["record"].source_completeness == "inspection_lot_only"]
    assert with_centre and lot_only

    output, filename = build_sample_register_workbook()
    assert filename.endswith(".xlsx")
    workbook = load_workbook(output)

    assert workbook.sheetnames[0] == "Summary"
    assert "Notification no lot" in workbook.sheetnames
    assert "Inspection lot only" in workbook.sheetnames

    def rows(name):
        sheet = workbook[name]
        return max(sheet.max_row - 2, 0)

    # A lot-only record has a lot by definition; that sheet must not be empty
    # just because the notification sheet next to it filters on having none.
    assert rows("Inspection lot only") == len(lot_only)

    subgroup_sheets = [
        name for name in workbook.sheetnames
        if name not in {"Summary", "Notification no lot", "Inspection lot only"}
    ]
    assert subgroup_sheets
    # Every work-centre record lands on exactly one sub-group sheet.
    assert sum(rows(name) for name in subgroup_sheets) == len(with_centre)


def test_register_export_is_not_capped_like_the_page(sap_app, monkeypatch):
    """The page truncates; the export must not inherit that."""
    from openpyxl import load_workbook
    from app.core.services import sap_quality_control
    from app.core.services.qc_register_export import build_sample_register_workbook
    from app.core.services.sap_quality_control import (
        import_sap_panvel_exports, sap_sample_register_data,
    )

    import_sap_panvel_exports(
        _inspection_export(), "SAP_INSPECTION_20260826.xlsx",
        _notification_export(), "SAP_NOTIFICATIONS_20260826.xlsx", None,
    )
    db.session.commit()
    total = sap_sample_register_data(limit=None)["total_matching"]

    monkeypatch.setattr(sap_quality_control, "SAP_REGISTER_VISIBLE_LIMIT", 1)
    assert len(sap_sample_register_data()["entries"]) == 1

    workbook = load_workbook(build_sample_register_workbook()[0])
    exported = sum(
        max(workbook[name].max_row - 2, 0)
        for name in workbook.sheetnames
        if name not in {"Summary", "Notification no lot", "Inspection lot only"}
    ) + max(workbook["Inspection lot only"].max_row - 2, 0)
    assert exported == total


def test_all_labs_presentation_groups_samples_by_lab_then_corporate_subgroup():
    from app.core.services.qc_presentation import _sap_presentation_action_groups

    panvel = {"code": "rgl_panvel", "name": "RGL Panvel"}
    chennai = {"code": "rgl_chennai", "name": "RGL Chennai"}

    def entry(laboratory, subgroup_key, subgroup_label, specification, record_id):
        return {
            "laboratory": laboratory,
            "subgroup_key": subgroup_key,
            "subgroup_label": subgroup_label,
            "specification_no": specification,
            "stt_overdue": False,
            "stt_due_date": date(2026, 8, 30),
            "record": SimpleNamespace(
                material_description=f"Sample {record_id}", notification_no=str(record_id), id=record_id,
            ),
        }

    groups = _sap_presentation_action_groups({
        "scope_laboratories": [panvel, chennai],
        "action_entries": [
            entry(panvel, "PC", "Production Chemicals", "ONGC / PC / 01 / 2026", 1),
            entry(chennai, "DFC", "Drilling Fluid Chemicals", "ONGC / DFC / 01 / 2026", 2),
            entry(panvel, "not_in_corporate_specification", "Not in Corporate Specification", None, 3),
            entry(panvel, "DFC", "Drilling Fluid Chemicals", "ONGC / DFC / 02 / 2026", 4),
        ],
    })

    assert [(group["laboratory"]["name"], group["subgroup_key"]) for group in groups] == [
        ("RGL Panvel", "DFC"),
        ("RGL Panvel", "PC"),
        ("RGL Panvel", "not_in_corporate_specification"),
        ("RGL Chennai", "DFC"),
    ]


def _presentation_readers():
    """A Corporate Chemistry reader and a laboratory one, for deck scoping."""
    from app.core.roles import SUPERUSER_ROLE, USER_ROLE
    from app.models.core.role import Role
    from app.models.core.user import User
    from app.models.core.user_module_permission import UserModulePermission

    user_role, super_role = Role(id=1, name=USER_ROLE), Role(id=2, name=SUPERUSER_ROLE)
    lab_user = User(
        id=1, username="panvel", password_hash="x", role=user_role, is_active=True,
        must_change_password=False, quality_control_lab_code="rgl_panvel",
    )
    superuser = User(
        id=2, username="super", password_hash="x", role=super_role,
        is_active=True, must_change_password=False,
    )
    db.session.add_all([
        user_role, super_role, lab_user, superuser,
        UserModulePermission(id=1, user_id=1, module_code="quality_control", can_access=True),
        UserModulePermission(id=2, user_id=2, module_code="quality_control", can_access=True),
    ])
    db.session.commit()
    return lab_user, superuser


def test_management_deck_is_built_at_the_reader_s_own_scope(sap_app, monkeypatch):
    """Corporate Chemistry chooses the scope; a laboratory always gets its own.

    The laboratory takes the same management deck to its review, narrowed to
    its own bench — and cannot widen it back by asking for another laboratory,
    or for the whole portfolio, in the query string.
    """
    from flask_login import login_user, logout_user
    from app.modules.quality_control.routes import download_portfolio_management_presentation

    captured = {}

    def build(static_folder, lab_codes):
        captured["lab_codes"] = lab_codes
        return BytesIO(b"presentation"), "presentation.pptx"

    monkeypatch.setattr(
        "app.core.services.qc_presentation.build_sap_portfolio_management_presentation", build,
    )
    lab_user, superuser = _presentation_readers()

    # Corporate Chemistry: a blank lab parameter still means every laboratory.
    with sap_app.test_request_context("/quality-control/management-review/presentation.pptx?lab="):
        login_user(superuser)
        response = download_portfolio_management_presentation.__wrapped__.__wrapped__()
        logout_user()
    assert response.status_code == 200
    assert captured["lab_codes"] is None

    # Corporate Chemistry: one named laboratory still narrows the deck.
    with sap_app.test_request_context(
        "/quality-control/management-review/presentation.pptx?lab=rgl_vadodara",
    ):
        login_user(superuser)
        download_portfolio_management_presentation.__wrapped__.__wrapped__()
        logout_user()
    assert captured["lab_codes"] == {"rgl_vadodara"}

    # A laboratory gets its own deck whatever the query string asks for.
    for query in ("", "?lab=", "?lab=rgl_vadodara", "?scope=labs"):
        with sap_app.test_request_context(
            f"/quality-control/management-review/presentation.pptx{query}",
        ):
            login_user(lab_user)
            response = download_portfolio_management_presentation.__wrapped__.__wrapped__()
            logout_user()
        assert response.status_code == 200, query
        assert captured["lab_codes"] == {"rgl_panvel"}, query


def test_open_register_paginates_every_item_for_the_management_presentation():
    from app.core.services.qc_presentation import _paginated_rows

    pages = _paginated_rows(list(range(37)), 11)

    assert [len(page) for page in pages] == [11, 11, 11, 4]
    assert [item for page in pages for item in page] == list(range(37))


def _year_inspection_export(rows: list[list[str]]) -> bytes:
    return _xlsx(pd.DataFrame(rows, columns=[
        "Inspection Lot", "Material", "Plant", "Start of Inspection", "End of Inspection",
        "System Status", "Usage Decision Code",
    ]))


def _year_notification_export(rows: list[list[str]], *, title: str) -> bytes:
    return _xlsx(pd.DataFrame(rows, columns=[
        "Notification No", "Notification Status", "Purchasing Document", "Item", "Material Number",
        "Material Description", "Work Center", "Plant", "Inspection Lot Number", "Status of Inspection Lot",
        "Start Date", "Planned End Date", "Completion Date", "Delay Days",
    ]), startrow=4, title=title)


def test_financial_year_is_named_from_the_first_of_april():
    from app.core.services.sap_quality_control import financial_year_label, financial_year_start

    assert financial_year_label(date(2026, 4, 1)) == "2026-27"
    assert financial_year_label(date(2026, 9, 1)) == "2026-27"
    assert financial_year_label(date(2027, 3, 31)) == "2026-27"
    assert financial_year_label(date(2027, 4, 1)) == "2027-28"
    assert financial_year_label(date(2099, 12, 1)) == "2099-00"
    assert financial_year_start(date(2027, 3, 31)) == date(2026, 4, 1)


def test_goods_receipt_lots_are_dropped_before_the_plant_routing_check():
    """A central QA33 selection returns lots from plants with no laboratory."""
    from app.core.services.sap_quality_control import parse_sap_inspection_workbook

    payload = parse_sap_inspection_workbook(
        _year_inspection_export([
            ["890000040001", "00001234", "10R2", "01.08.2026", "", "REL CALC", ""],
            # A goods-receipt lot raised at a plant outside the routing map.
            ["100001466468", "00009999", "12F1", "01.08.2026", "", "INSP RREC", ""],
        ]),
        "SAP_INSP_20260901.xlsx", expected_plant=None, allow_multiple_plants=True,
    )

    assert [row["inspection_lot_number"] for row in payload.rows] == ["890000040001"]
    assert payload.excluded_rows == {"non_laboratory_lots": 1}


def test_notifications_before_the_financial_year_are_left_out_of_the_base():
    from app.core.services.sap_quality_control import parse_sap_notification_workbook

    payload = parse_sap_notification_workbook(
        _year_notification_export([
            ["020000030000", "COMP", "45000001", "10", "00001234", "Barytes", "MUDLAB", "10R2",
             "890000030000", "UD ICCO", "12.11.2025", "20.11.2025", "25.11.2025", "1"],
            ["020000036313", "OPEN", "45000002", "10", "00005678", "Glycol", "OILLAB", "10R2",
             "890000041007", "REL CALC", "01.04.2026", "10.04.2026", "", "0"],
            # SAP dates neither the raising nor the close of this notification.
            ["020000000001", "COMP", "45000003", "10", "00007777", "Undated", "OILLAB", "10R2",
             "890000000001", "UD ICCO", "", "", "", ""],
        ], title="Date : 01.09.2026"),
        "SAP_ZLABIMS_20260901.xlsx", expected_plant="10R2",
    )

    assert [row["notification_no"] for row in payload.rows] == ["020000036313"]
    assert payload.excluded_rows == {"before_financial_year": 1, "no_start_date": 1}


def test_a_daily_upload_of_current_work_leaves_the_years_base_data_in_view(sap_app):
    """The base is loaded once; each morning's export carries far less.

    Every screen reads the financial year rather than the newest batch, so a
    notification the base recorded stays in the register and on the dashboard
    even on a day SAP no longer reports it.
    """
    from app.core.services.sap_quality_control import (
        import_sap_lab_exports, sap_lab_dashboard_data, sap_sample_register_data,
    )
    from app.models.quality_control.qc_sap_monitoring import QCSAPRecord

    import_sap_lab_exports(
        "rgl_panvel",
        _year_inspection_export([
            ["890000040001", "00001234", "10R2", "05.04.2026", "20.04.2026", "UD ICCO STUP", "A"],
            ["890000040002", "00005678", "10R2", "06.04.2026", "", "REL CALC", ""],
        ]),
        "SAP_INSP_20260430.xlsx",
        _year_notification_export([
            ["020000040001", "COMP", "45000001", "10", "00001234", "Barytes", "MUDLAB", "10R2",
             "890000040001", "UD ICCO", "05.04.2026", "18.04.2026", "20.04.2026", "2"],
            ["020000040002", "OPEN", "45000002", "10", "00005678", "Glycol", "OILLAB", "10R2",
             "890000040002", "REL CALC", "06.04.2026", "20.04.2026", "", "1"],
        ], title="Date : 30.04.2026"),
        "SAP_ZLABIMS_20260430.xlsx", None,
    )
    db.session.commit()

    # The next morning SAP reports only what is still current, plus one new
    # notification raised that day.
    import_sap_lab_exports(
        "rgl_panvel",
        _year_inspection_export([
            ["890000040002", "00005678", "10R2", "06.04.2026", "", "REL CALC", ""],
            ["890000040003", "00009999", "10R2", "01.09.2026", "", "CRTD CHCR", ""],
        ]),
        "SAP_INSP_20260901.xlsx",
        _year_notification_export([
            ["020000040002", "OPEN", "45000002", "10", "00005678", "Glycol", "OILLAB", "10R2",
             "890000040002", "REL CALC", "06.04.2026", "20.04.2026", "", "1"],
            ["020000040003", "OPEN", "45000003", "10", "00009999", "Xylene", "QUALILAB", "10R2",
             "890000040003", "CRTD CHCR", "01.09.2026", "10.09.2026", "", "0"],
        ], title="Date : 01.09.2026"),
        "SAP_ZLABIMS_20260901.xlsx", None,
    )
    db.session.commit()

    assert {record.financial_year for record in QCSAPRecord.query.all()} == {"2026-27"}

    dashboard = sap_lab_dashboard_data("rgl_panvel")
    on_screen = {entry["record"].notification_no for entry in dashboard["records"]}
    assert on_screen == {"020000040001", "020000040002", "020000040003"}
    assert dashboard["batch"].as_of_date == date(2026, 9, 1)

    register = sap_sample_register_data(lab_code="rgl_panvel")
    assert {entry["record"].notification_no for entry in register["entries"]} == on_screen


def test_a_closed_year_is_not_mixed_into_the_year_being_reported(sap_app):
    from app.core.services.sap_quality_control import import_sap_lab_exports, sap_lab_dashboard_data
    from app.models.quality_control.qc_sap_monitoring import QCSAPRecord

    import_sap_lab_exports(
        "rgl_panvel",
        _year_inspection_export([
            ["890000030001", "00001234", "10R2", "05.09.2025", "20.09.2025", "UD ICCO STUP", "A"],
        ]),
        "SAP_INSP_20260131.xlsx",
        _year_notification_export([
            ["020000030001", "COMP", "45000001", "10", "00001234", "Barytes", "MUDLAB", "10R2",
             "890000030001", "UD ICCO", "05.09.2025", "18.09.2025", "20.09.2025", "2"],
        ], title="Date : 31.01.2026"),
        "SAP_ZLABIMS_20260131.xlsx", None,
    )
    import_sap_lab_exports(
        "rgl_panvel",
        _year_inspection_export([
            ["890000040001", "00005678", "10R2", "05.04.2026", "", "REL CALC", ""],
        ]),
        "SAP_INSP_20260901.xlsx",
        _year_notification_export([
            ["020000040001", "OPEN", "45000002", "10", "00005678", "Glycol", "OILLAB", "10R2",
             "890000040001", "REL CALC", "05.04.2026", "20.04.2026", "", "1"],
        ], title="Date : 01.09.2026"),
        "SAP_ZLABIMS_20260901.xlsx", None,
    )
    db.session.commit()

    assert {
        (record.notification_no, record.financial_year) for record in QCSAPRecord.query.all()
    } == {("020000030001", "2025-26"), ("020000040001", "2026-27")}

    dashboard = sap_lab_dashboard_data("rgl_panvel")
    assert [entry["record"].notification_no for entry in dashboard["records"]] == ["020000040001"]


def test_financial_year_scope_states_the_span_and_rolls_over_on_its_own():
    """The wording is derived from the data, so April re-dates it unattended."""
    from app.core.services.sap_quality_control import financial_year_scope

    scope = financial_year_scope(date(2026, 9, 1))
    assert scope["label"] == "2026-27"
    assert scope["start_date"] == date(2026, 4, 1)
    assert scope["note"] == "All SAP notifications created on or after 01.04.2026"

    # The last day of the year still names that year, and the first day of the
    # next names the next -- no edit is needed for the wording to stay true.
    assert financial_year_scope(date(2027, 3, 31))["note"].endswith("01.04.2026")
    assert financial_year_scope(date(2027, 4, 1))["note"].endswith("01.04.2027")


def test_every_sap_screen_states_the_financial_year_it_is_reporting(sap_app):
    from app.core.services.sap_quality_control import (
        import_sap_lab_exports, sap_lab_dashboard_data, sap_management_data,
        sap_sample_register_data,
    )

    import_sap_lab_exports(
        "rgl_panvel",
        _year_inspection_export([
            ["890000040001", "00001234", "10R2", "05.04.2026", "", "REL CALC", ""],
        ]),
        "SAP_INSP_20260901.xlsx",
        _year_notification_export([
            ["020000040001", "OPEN", "45000001", "10", "00001234", "Barytes", "MUDLAB", "10R2",
             "890000040001", "REL CALC", "05.04.2026", "20.04.2026", "", "1"],
        ], title="Date : 01.09.2026"),
        "SAP_ZLABIMS_20260901.xlsx", None,
    )
    db.session.commit()

    expected = {
        "label": "2026-27",
        "start_date": date(2026, 4, 1),
        "note": "All SAP notifications created on or after 01.04.2026",
    }
    assert sap_lab_dashboard_data("rgl_panvel")["financial_year_scope"] == expected
    assert sap_sample_register_data(lab_code="rgl_panvel")["financial_year_scope"] == expected
    assert sap_management_data()["financial_year_scope"] == expected


def test_a_laboratory_awaiting_its_first_import_claims_no_financial_year(sap_app):
    from app.core.services.sap_quality_control import sap_lab_dashboard_data

    assert sap_lab_dashboard_data("rgl_panvel")["financial_year_scope"] is None


def test_a_year_rebuild_replaces_sap_fields_without_touching_the_human_layer(sap_app):
    """A rebuild is a refresh of SAP's facts, not a wipe.

    The laboratory's returned follow-up and the QC-admin exclusion are the two
    things SAP does not hold and cannot reproduce, so a rebuild that dropped
    them would destroy the only copy.
    """
    from app.core.services.sap_quality_control import (
        create_sap_lab_update, exclude_sap_record_from_monitoring,
        import_sap_lab_exports, rebuild_sap_financial_year,
    )
    from app.models.quality_control.qc_sap_monitoring import QCSAPLabUpdate, QCSAPRecord

    import_sap_lab_exports(
        "rgl_panvel",
        _year_inspection_export([
            ["890000040001", "00001234", "10R2", "05.04.2026", "", "REL CALC", ""],
            ["890000040002", "00005678", "10R2", "06.04.2026", "", "REL CALC", ""],
        ]),
        "SAP_INSP_20260430.xlsx",
        _year_notification_export([
            ["020000040001", "OPEN", "45000001", "10", "00001234", "Barytes", "MUDLAB", "10R2",
             "890000040001", "REL CALC", "05.04.2026", "20.04.2026", "", "1"],
            ["020000040002", "OPEN", "45000002", "10", "00005678", "Glycol", "OILLAB", "10R2",
             "890000040002", "REL CALC", "06.04.2026", "20.04.2026", "", "1"],
        ], title="Date : 30.04.2026"),
        "SAP_ZLABIMS_20260430.xlsx", None,
    )
    db.session.commit()

    followed = QCSAPRecord.query.filter_by(notification_no="020000040001").one()
    excluded = QCSAPRecord.query.filter_by(notification_no="020000040002").one()
    create_sap_lab_update(
        followed.id,
        {"activity_status": "under_testing", "action_owner": "Bench 2", "update_note": "Retest running"},
        None, lab_code="rgl_panvel",
    )
    exclude_sap_record_from_monitoring(
        excluded.id, {"exclusion_reason": "junk_notification", "note": "Raised in error"},
        None, lab_code="rgl_panvel",
    )
    db.session.commit()

    # The rebuild sees the first lot closed and decided, and no longer carries
    # the second notification at all.
    result = rebuild_sap_financial_year(
        _year_inspection_export([
            ["890000040001", "00001234", "10R2", "05.04.2026", "28.04.2026", "UD ICCO STUP", "A"],
        ]),
        "SAP_INSP_20260831.xlsx",
        _year_notification_export([
            ["020000040001", "COMP", "45000001", "10", "00001234", "Barytes", "MUDLAB", "10R2",
             "890000040001", "UD ICCO", "05.04.2026", "20.04.2026", "28.04.2026", "1"],
        ], title="Date : 01.09.2026"),
        "SAP_ZLABIMS_20260901.xlsx", None, as_of_date=date(2026, 9, 1),
    )
    db.session.commit()

    assert result["financial_year"] == "2026-27"
    # SAP's own fields were replaced in place on the record that remains.
    refreshed = QCSAPRecord.query.filter_by(notification_no="020000040001").one()
    assert refreshed.id == followed.id
    assert refreshed.official_status == "completed"
    assert refreshed.usage_decision_code == "A"
    assert QCSAPLabUpdate.query.filter_by(record_id=refreshed.id).count() == 1

    # The excluded record is absent from the exports, but its QC-admin decision
    # keeps it: it is reported for review rather than deleted.
    reconciliation = result["reconciliation"]
    assert reconciliation["removed_count"] == 0
    assert reconciliation["retained_count"] == 1
    assert reconciliation["retained"][0]["notification_no"] == "020000040002"
    assert reconciliation["retained"][0]["reason"] == "QC-admin monitoring decision recorded"
    assert QCSAPRecord.query.filter_by(notification_no="020000040002").count() == 1


def test_a_year_rebuild_retires_only_records_nothing_is_recorded_against(sap_app):
    from app.core.services.sap_quality_control import (
        import_sap_lab_exports, rebuild_sap_financial_year,
    )
    from app.models.quality_control.qc_sap_monitoring import QCSAPRecord

    import_sap_lab_exports(
        "rgl_panvel",
        _year_inspection_export([
            ["890000040001", "00001234", "10R2", "05.04.2026", "", "REL CALC", ""],
            ["890000040009", "00009999", "10R2", "07.04.2026", "", "REL CALC", ""],
        ]),
        "SAP_INSP_20260430.xlsx",
        _year_notification_export([
            ["020000040001", "OPEN", "45000001", "10", "00001234", "Barytes", "MUDLAB", "10R2",
             "890000040001", "REL CALC", "05.04.2026", "20.04.2026", "", "1"],
            ["020000040009", "OPEN", "45000009", "10", "00009999", "Withdrawn", "OILLAB", "10R2",
             "890000040009", "REL CALC", "07.04.2026", "20.04.2026", "", "1"],
        ], title="Date : 30.04.2026"),
        "SAP_ZLABIMS_20260430.xlsx", None,
    )
    db.session.commit()
    assert QCSAPRecord.query.count() == 2

    result = rebuild_sap_financial_year(
        _year_inspection_export([
            ["890000040001", "00001234", "10R2", "05.04.2026", "", "REL CALC", ""],
        ]),
        "SAP_INSP_20260831.xlsx",
        _year_notification_export([
            ["020000040001", "OPEN", "45000001", "10", "00001234", "Barytes", "MUDLAB", "10R2",
             "890000040001", "REL CALC", "05.04.2026", "20.04.2026", "", "1"],
        ], title="Date : 01.09.2026"),
        "SAP_ZLABIMS_20260901.xlsx", None, as_of_date=date(2026, 9, 1),
    )
    db.session.commit()

    assert result["reconciliation"]["removed_count"] == 1
    assert result["reconciliation"]["removed_sample"] == ["020000040009"]
    assert [record.notification_no for record in QCSAPRecord.query.all()] == ["020000040001"]

    # Left off, the same rebuild keeps everything the exports omit.
    assert rebuild_sap_financial_year(
        _year_inspection_export([
            ["890000040001", "00001234", "10R2", "05.04.2026", "", "REL CALC", ""],
        ]),
        "SAP_INSP_20260831.xlsx",
        _year_notification_export([
            ["020000040001", "OPEN", "45000001", "10", "00001234", "Barytes", "MUDLAB", "10R2",
             "890000040001", "REL CALC", "05.04.2026", "20.04.2026", "", "1"],
        ], title="Date : 01.09.2026"),
        "SAP_ZLABIMS_20260901.xlsx", None, as_of_date=date(2026, 9, 1), reconcile=False,
    )["reconciliation"] is None


def test_each_import_records_what_it_moved_since_the_records_already_held(sap_app):
    import json

    from app.core.services.sap_quality_control import import_sap_lab_exports

    def upload(lot_status, usage, notif_status, completion, title, filename):
        return import_sap_lab_exports(
            "rgl_panvel",
            _year_inspection_export([
                ["890000040001", "00001234", "10R2", "05.04.2026", "", lot_status, usage],
                ["890000040002", "00005678", "10R2", "06.04.2026", "", "REL CALC", ""],
            ]),
            filename,
            _year_notification_export([
                ["020000040001", notif_status, "45000001", "10", "00001234", "Barytes", "MUDLAB",
                 "10R2", "890000040001", lot_status, "05.04.2026", "20.04.2026", completion, "1"],
                ["020000040002", "OPEN", "45000002", "10", "00005678", "Glycol", "OILLAB",
                 "10R2", "890000040002", "REL CALC", "06.04.2026", "20.04.2026", "", "1"],
            ], title=title),
            filename.replace("INSP", "ZLABIMS"), None,
        )

    first = upload("REL CALC", "", "OPEN", "", "Date : 30.04.2026", "SAP_INSP_20260430.xlsx")
    db.session.commit()
    changes = json.loads(first.summary_json)["changes"]
    assert changes["new"]["count"] == 2
    assert changes["closed"]["count"] == 0

    second = upload("UD ICCO STUP", "A", "COMP", "28.04.2026", "Date : 01.09.2026", "SAP_INSP_20260901.xlsx")
    db.session.commit()
    changes = json.loads(second.summary_json)["changes"]
    assert changes["new"]["count"] == 0
    assert changes["closed"] == {"count": 1, "sample": ["020000040001"]}
    assert changes["usage_decided"] == {"count": 1, "sample": ["020000040001"]}
    assert changes["reopened"]["count"] == 0


def test_a_central_export_sets_aside_plants_no_laboratory_owns(sap_app):
    """A company-wide QA33 selection is not a broken export.

    SAP raises laboratory lots at plants with no RGL or IDWE of their own.
    Refusing the whole upload over them blocked the day's monitoring, so they
    are set aside and counted instead.
    """
    from app.core.services.sap_quality_control import import_central_sap_exports
    from app.models.quality_control.qc_sap_monitoring import QCSAPRecord

    batches = import_central_sap_exports(
        _year_inspection_export([
            ["890000040001", "00001234", "10R2", "05.04.2026", "", "REL CALC", ""],
            ["890000040002", "00005678", "42R2", "06.04.2026", "", "REL CALC", ""],
            # Laboratory lots raised at plants outside the approved routing.
            ["890000040003", "00009999", "22A1", "06.04.2026", "", "REL CALC", ""],
            ["890000040004", "00009999", "20A1", "06.04.2026", "", "REL CALC", ""],
        ]),
        "SAP_INSP_20260901.xlsx",
        _year_notification_export([
            ["020000040001", "OPEN", "45000001", "10", "00001234", "Barytes", "MUDLAB", "10R2",
             "890000040001", "REL CALC", "05.04.2026", "20.04.2026", "", "1"],
            ["020000040002", "OPEN", "45000002", "10", "00005678", "Glycol", "QUALILAB", "42R2",
             "890000040002", "REL CALC", "06.04.2026", "20.04.2026", "", "1"],
        ], title="Date : 01.09.2026"),
        "SAP_ZLABIMS_20260901.xlsx", None,
    )
    db.session.commit()

    assert {batch.plant_code for batch in batches} == {"10R2", "42R2"}
    assert json.loads(batches[0].summary_json)["excluded_rows"]["lots_outside_rgl_idwe"] == 2
    assert QCSAPRecord.query.count() == 2
    assert not QCSAPRecord.query.filter(
        QCSAPRecord.inspection_lot_number.in_(["890000040003", "890000040004"])
    ).count()


def test_an_export_holding_no_laboratory_plant_at_all_still_stops_the_upload(sap_app):
    """Setting aside foreign plants must not swallow a wholly wrong report."""
    import pytest

    from app.core.services.sap_quality_control import import_central_sap_exports

    with pytest.raises(ValueError, match="no rows for an RGL or IDWE plant"):
        import_central_sap_exports(
            _year_inspection_export([
                ["890000040003", "00009999", "22A1", "06.04.2026", "", "REL CALC", ""],
            ]),
            "SAP_INSP_20260901.xlsx",
            _year_notification_export([
                ["020000040001", "OPEN", "45000001", "10", "00001234", "Barytes", "MUDLAB", "10R2",
                 "890000040001", "REL CALC", "05.04.2026", "20.04.2026", "", "1"],
            ], title="Date : 01.09.2026"),
            "SAP_ZLABIMS_20260901.xlsx", None,
        )


def test_one_upload_stores_its_workbooks_once_for_every_laboratory(sap_app):
    """A central upload writes a batch per laboratory from a single pair."""
    from app.core.services.sap_quality_control import import_central_sap_exports
    from app.models.quality_control.qc_sap_monitoring import QCSAPSourceDocument

    batches = import_central_sap_exports(
        _central_inspection_export(), "SAP_INSP_20260827.xlsx",
        _central_notification_export(), "SAP_ZLABIMS_20260827.xlsx", None,
    )
    db.session.commit()

    assert len(batches) == 3
    assert len({batch.source_document_id for batch in batches}) == 1
    assert QCSAPSourceDocument.query.count() == 1
    document = QCSAPSourceDocument.query.one()
    assert document.purged_at is None
    assert document.inspection_source_data
    # Each batch still describes what it was built from.
    assert all(batch.inspection_filename == "SAP_INSP_20260827.xlsx" for batch in batches)
    assert all(batch.source_is_available for batch in batches)


def test_a_new_upload_supersedes_the_previous_workbooks(sap_app):
    """Only the current pair is retained; the one before it is cleared."""
    from app.core.services.sap_quality_control import import_central_sap_exports
    from app.models.quality_control.qc_sap_monitoring import QCSAPSourceDocument, QCSAPUploadBatch

    first = import_central_sap_exports(
        _central_inspection_export(), "SAP_INSP_20260827.xlsx",
        _central_notification_export(), "SAP_ZLABIMS_20260827.xlsx", None,
    )
    db.session.commit()
    first_document_id = first[0].source_document_id

    # The same day re-uploaded: still a new pair, still supersedes the old one.
    import_central_sap_exports(
        _central_inspection_export(), "SAP_INSP_20260827.xlsx",
        _central_notification_export(), "SAP_ZLABIMS_20260827.xlsx", None,
    )
    db.session.commit()

    superseded = db.session.get(QCSAPSourceDocument, first_document_id)
    assert superseded.purged_at is not None
    assert superseded.inspection_source_data == b""
    assert superseded.notification_source_data == b""

    current = QCSAPSourceDocument.query.filter(
        QCSAPSourceDocument.purged_at.is_(None)
    ).one()
    assert current.inspection_source_data

    # The superseded upload's batches survive with their metadata; only the
    # workbooks are gone, and the dashboard can say so.
    old_batches = QCSAPUploadBatch.query.filter_by(source_document_id=first_document_id).all()
    assert old_batches
    assert all(not batch.source_is_available for batch in old_batches)
    assert all(batch.record_count for batch in old_batches)


def test_an_older_daily_export_cannot_regress_the_current_sap_position(sap_app):
    """The live batch label and canonical record values must never disagree."""
    from app.core.services.sap_quality_control import import_sap_lab_exports
    from app.models.quality_control.qc_sap_monitoring import QCSAPRecord, QCSAPSourceDocument

    current_inspection = _year_inspection_export([
        ["890000040001", "00001234", "10R2", "01.09.2026", "", "REL CALC", ""],
    ])
    current_notifications = _year_notification_export([
        ["020000040001", "OPEN", "45000001", "10", "00001234", "Barytes", "NEWLAB", "10R2",
         "890000040001", "REL CALC", "01.09.2026", "10.09.2026", "", "0"],
    ], title="Date : 02.09.2026")
    import_sap_lab_exports(
        "rgl_panvel", current_inspection, "SAP_INSP_20260902.xlsx",
        current_notifications, "SAP_ZLABIMS_20260902.xlsx", None,
    )
    db.session.commit()

    with pytest.raises(ValueError, match="older than the current SAP snapshot"):
        import_sap_lab_exports(
            "rgl_panvel",
            _year_inspection_export([
                ["890000040001", "00001234", "10R2", "01.09.2026", "", "REL CALC", ""],
            ]),
            "SAP_INSP_20260901.xlsx",
            _year_notification_export([
                ["020000040001", "OPEN", "45000001", "10", "00001234", "Barytes", "OLDLAB", "10R2",
                 "890000040001", "REL CALC", "01.09.2026", "10.09.2026", "", "0"],
            ], title="Date : 01.09.2026"),
            "SAP_ZLABIMS_20260901.xlsx", None,
        )

    record = QCSAPRecord.query.filter_by(notification_no="020000040001").one()
    assert record.work_center == "NEWLAB"
    assert QCSAPSourceDocument.query.count() == 1
    assert QCSAPSourceDocument.query.one().purged_at is None


def test_a_single_lab_upload_keeps_other_labs_current_source_pair(sap_app):
    """A compatible one-lab upload must not erase the central pair it leaves live."""
    from app.core.services.sap_quality_control import (
        import_central_sap_exports, import_sap_lab_exports, latest_sap_batch,
    )
    from app.models.quality_control.qc_sap_monitoring import QCSAPSourceDocument

    central_batches = import_central_sap_exports(
        _central_inspection_export(), "SAP_INSP_20260827.xlsx",
        _central_notification_export(), "SAP_ZLABIMS_20260827.xlsx", None,
    )
    db.session.commit()
    central_document_id = central_batches[0].source_document_id

    import_sap_lab_exports(
        "rgl_panvel",
        _year_inspection_export([
            ["890000040001", "00001234", "10R2", "01.09.2026", "", "REL CALC", ""],
        ]),
        "SAP_INSP_20260901.xlsx",
        _year_notification_export([
            ["020000040001", "OPEN", "45000001", "10", "00001234", "Barytes", "NEWLAB", "10R2",
             "890000040001", "REL CALC", "01.09.2026", "10.09.2026", "", "0"],
        ], title="Date : 01.09.2026"),
        "SAP_ZLABIMS_20260901.xlsx", None,
    )
    db.session.commit()

    panvel_document_id = latest_sap_batch("rgl_panvel").source_document_id
    chennai_batch = latest_sap_batch("rgl_chennai")
    retained = {
        document.id for document in QCSAPSourceDocument.query.filter(
            QCSAPSourceDocument.purged_at.is_(None),
        ).all()
    }
    assert retained == {central_document_id, panvel_document_id}
    assert chennai_batch.source_document_id == central_document_id
    assert chennai_batch.source_is_available is True


def test_the_retention_sweep_repairs_a_purge_missed_at_import(sap_app):
    from app.core.services.audit_workbook_retention import purge_expired_audit_workbook_payloads
    from app.core.services.sap_quality_control import import_central_sap_exports
    from app.models.quality_control.qc_sap_monitoring import QCSAPSourceDocument

    import_central_sap_exports(
        _central_inspection_export(), "SAP_INSP_20260827.xlsx",
        _central_notification_export(), "SAP_ZLABIMS_20260827.xlsx", None,
    )
    db.session.commit()
    # A second pair left behind as though its supersede had not run.
    stray = QCSAPSourceDocument(
        inspection_filename="stray.xlsx", inspection_content_type="x", inspection_file_size=3,
        inspection_source_data=b"abc",
        notification_filename="stray2.xlsx", notification_content_type="x", notification_file_size=3,
        notification_source_data=b"abc",
        uploaded_at=datetime(2026, 8, 1, 9, 0),
    )
    db.session.add(stray)
    db.session.commit()

    counts = purge_expired_audit_workbook_payloads()
    db.session.commit()

    assert counts["qc_sap"] == 1
    assert db.session.get(QCSAPSourceDocument, stray.id).purged_at is not None
    assert QCSAPSourceDocument.query.filter(QCSAPSourceDocument.purged_at.is_(None)).count() == 1


def test_the_source_audit_trail_lists_only_the_most_recent_uploads(sap_app):
    """The trail is trimmed, and the page says how far back it goes."""
    from app.core.services.sap_quality_control import (
        SAP_SOURCE_AUDIT_TRAIL_LIMIT, import_sap_lab_exports, sap_lab_dashboard_data,
    )
    from app.models.quality_control.qc_sap_monitoring import QCSAPUploadBatch

    import_sap_lab_exports(
        "rgl_panvel",
        _year_inspection_export([
            ["890000040001", "00001234", "10R2", "05.04.2026", "", "REL CALC", ""],
        ]),
        "SAP_INSP_20260901.xlsx",
        _year_notification_export([
            ["020000040001", "OPEN", "45000001", "10", "00001234", "Barytes", "MUDLAB", "10R2",
             "890000040001", "REL CALC", "05.04.2026", "20.04.2026", "", "1"],
        ], title="Date : 01.09.2026"),
        "SAP_ZLABIMS_20260901.xlsx", None,
    )
    # The rest are recorded directly: this is about how many the page lists,
    # not about parsing the same workbook a dozen more times.
    for day in range(2, SAP_SOURCE_AUDIT_TRAIL_LIMIT + 4):
        db.session.add(QCSAPUploadBatch(
            lab_code="rgl_panvel", plant_code="10R2", as_of_date=date(2026, 9, day),
            inspection_filename="i.xlsx", inspection_content_type="x", inspection_file_size=1,
            notification_filename="n.xlsx", notification_content_type="x", notification_file_size=1,
            record_count=1, summary_json="{}",
        ))
    db.session.commit()

    total = QCSAPUploadBatch.query.filter_by(lab_code="rgl_panvel").count()
    assert total == SAP_SOURCE_AUDIT_TRAIL_LIMIT + 3

    data = sap_lab_dashboard_data("rgl_panvel")
    assert data["source_audit_trail_limit"] == SAP_SOURCE_AUDIT_TRAIL_LIMIT
    assert len(data["recent_batches"]) == SAP_SOURCE_AUDIT_TRAIL_LIMIT
    # Newest first, and the trimmed-away uploads are the oldest ones.
    listed = [batch.as_of_date for batch in data["recent_batches"]]
    assert listed == sorted(listed, reverse=True)
    assert listed[0] == date(2026, 9, SAP_SOURCE_AUDIT_TRAIL_LIMIT + 3)
    assert date(2026, 9, 1) not in listed
