"""On-demand PowerPoint exports of the Inventory Monitoring management reviews."""
from __future__ import annotations

from datetime import date
from decimal import Decimal
from io import BytesIO
from pathlib import Path
from typing import Any

INK, TEAL, RED, AMBER, PURPLE, GREY, BORDER = "0F3A44", "0B8278", "B3261E", "B45309", "6D28D9", "5F7379", "D8E8E5"
BAND_TONES = {
    "critical_low_stock": RED, "low_stock": AMBER, "healthy_stock": TEAL,
    "slow_moving_stock": AMBER, "excess_stock": PURPLE, "unclassified": GREY,
}
# The unit follows the material code in every register, as it does on the pages.
REGISTER_HEADERS = ["Material", "UoM", "Description", "Grp", "Work centre", "Zone", "Stock qty", "Value", "Months"]
REGISTER_WIDTHS = [1.35, .55, 2.75, .5, 2.4, 1.6, 1.2, 1.3, .8]
SUPPORTING_HEADERS = ["Material", "UoM", "Description", "Grp", "Work centre", "Zone", "Value"]
SUPPORTING_WIDTHS = [1.4, .55, 3.45, .6, 2.6, 1.85, 2.0]
CENTRE_REGISTER_HEADERS = ["Material", "UoM", "Description", "Grp", "Stock qty", "Inventory value", "Months of cover"]
CENTRE_REGISTER_WIDTHS = [1.5, .55, 4.6, .6, 1.7, 1.9, 1.6]
KICKER = "ONGC CORPORATE CHEMISTRY · INVENTORY MONITORING · GROUPS 09 AND 10"


def _crore(value: Any) -> str:
    return f"₹ {float(value or 0) / 10000000:,.2f} Cr"


def _months(value: Any) -> str:
    return "—" if value is None else f"{float(value):,.1f}"


def _quantity(value: Any, uom: str | None) -> str:
    return "—" if value is None else f"{float(value):,.0f} {uom or ''}".strip()


def _concise(value: Any, limit: int = 58) -> str:
    text = " ".join(str(value or "—").split())
    return text if len(text) <= limit else f"{text[:limit + 1].rsplit(' ', 1)[0].rstrip('.,;:')}…"


class _Deck:
    """Shared 16:9 chrome for the Inventory Monitoring decks, loaded only when an export is requested."""

    def __init__(self, static_folder: str, source_line: str) -> None:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        self._rgb, self._auto_shape, self.inches, self.pt = RGBColor, MSO_SHAPE, Inches, Pt
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = Inches(13.333), Inches(7.5)
        self.blank = self.prs.slide_layouts[6]
        self.source_line = source_line
        self.corporate_chemistry_logo = Path(static_folder) / "images" / "ongc-corporate-chemistry-logo.png"
        self.ongc_logo = Path(static_folder) / "images" / "ongc-official-logo.png"

    def color(self, value: str) -> Any:
        return self._rgb.from_string(value)

    def add_text(self, slide, value, x, y, w, h, size=18, tone=INK, bold=False, wrap=False):
        shape = slide.shapes.add_textbox(self.inches(x), self.inches(y), self.inches(w), self.inches(h))
        shape.text_frame.word_wrap = wrap
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.text = str(value)
        paragraph.font.size, paragraph.font.bold, paragraph.font.name = self.pt(size), bold, "Arial"
        paragraph.font.color.rgb = self.color(tone)
        return shape

    def rectangle(self, slide, x, y, w, h, fill, line=None):
        shape = slide.shapes.add_shape(self._auto_shape.RECTANGLE, self.inches(x), self.inches(y), self.inches(w), self.inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color(fill)
        shape.line.color.rgb = self.color(line or fill)
        return shape

    def background(self, slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.color("FFFFFF")
        self.rectangle(slide, 0, 0, 13.333, .08, INK)
        self.rectangle(slide, .42, 1.26, 12.45, .015, BORDER)
        self.rectangle(slide, .42, 6.93, 12.45, .015, BORDER)

    def add_header_branding(self, slide):
        """Keep the ONGC and Corporate Chemistry marks together in every header."""
        if self.ongc_logo.exists():
            slide.shapes.add_picture(str(self.ongc_logo), self.inches(11.15), self.inches(.21), width=self.inches(.9), height=self.inches(.51))
        if self.corporate_chemistry_logo.exists():
            slide.shapes.add_picture(str(self.corporate_chemistry_logo), self.inches(12.24), self.inches(.16), width=self.inches(.55), height=self.inches(.55))

    def add_cover_branding(self, slide):
        self.rectangle(slide, 10.0, .46, 2.32, 1.2, "FFFFFF", BORDER)
        if self.ongc_logo.exists():
            slide.shapes.add_picture(str(self.ongc_logo), self.inches(10.16), self.inches(.68), width=self.inches(.9), height=self.inches(.51))
        if self.corporate_chemistry_logo.exists():
            slide.shapes.add_picture(str(self.corporate_chemistry_logo), self.inches(11.16), self.inches(.53), width=self.inches(1.05), height=self.inches(1.05))

    def header(self, slide, title):
        self.background(slide)
        self.add_text(slide, KICKER, .42, .27, 8.6, .24, 11, TEAL, True)
        self.add_text(slide, title, .42, .7, 11.35, .46, 26, INK, True)
        self.add_header_branding(slide)
        self.add_text(slide, self.source_line, .42, 7.08, 9.6, .16, 8, GREY)
        self.add_text(slide, f"{len(self.prs.slides._sldIdLst):02d}", 12.5, 7.08, .25, .16, 8, GREY)

    def new_slide(self, title):
        slide = self.prs.slides.add_slide(self.blank)
        self.header(slide, title)
        return slide

    def cover(self, kicker, title, period, subtitle):
        slide = self.prs.slides.add_slide(self.blank)
        self.background(slide)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.color("E7F5F2")
        self.rectangle(slide, 12.48, .08, .85, 6.85, "D2EBE6")
        self.add_text(slide, kicker, .76, 1.28, 8.0, .28, 14, TEAL, True)
        self.add_text(slide, title, .76, 1.82, 10.6, .8, 44, INK, True, wrap=True)
        self.rectangle(slide, .76, 2.82, 1.15, .045, TEAL)
        self.add_text(slide, period, .76, 3.14, 8.5, .36, 24, INK, True)
        self.add_text(slide, subtitle, .76, 3.74, 10.4, .6, 16, GREY, wrap=True)
        self.add_cover_branding(slide)
        self.add_text(slide, "Office of Head Corporate Chemistry | Mumbai / Dehradun", .76, 6.45, 7.2, .2, 11, GREY)
        return slide

    def metric(self, slide, x, y, value, label, tone=INK, note=None):
        card = slide.shapes.add_shape(self._auto_shape.ROUNDED_RECTANGLE, self.inches(x - .18), self.inches(y - .18), self.inches(3.45), self.inches(1.35 if note is None else 1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = self.color({TEAL: "E7F5F2", RED: "FBEBEA", AMBER: "FDF1E3", PURPLE: "F1EBFB"}.get(tone, "F1F5F5"))
        card.line.color.rgb = self.color(BORDER)
        self.add_text(slide, value, x, y, 3.0, .4, 26, tone, True)
        self.add_text(slide, label, x, y + .52, 3.05, .25, 13, INK, True)
        if note:
            self.add_text(slide, note, x, y + .82, 3.05, .25, 10, GREY, wrap=True)

    def table(self, slide, headers, rows, widths, y=1.45, font_size=8, empty_note="No lines fall in this register for the selected reporting date."):
        if not rows:
            self.add_text(slide, empty_note, .8, 2.1, 11.0, .35, 17, TEAL, True, wrap=True)
            return
        height = min(5.25, .3 * (len(rows) + 1))
        shape = slide.shapes.add_table(len(rows) + 1, len(headers), self.inches(.42), self.inches(y), self.inches(12.45), self.inches(height)).table
        for index, width in enumerate(widths):
            shape.columns[index].width = self.inches(width)
        for column, label in enumerate(headers):
            cell = shape.cell(0, column)
            cell.text = label
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.color(INK)
        sections = []
        for row_index, values in enumerate(rows, 1):
            if isinstance(values, dict):  # a specification sub-heading spanning the table
                sections.append(row_index)
                cell = shape.cell(row_index, 0)
                cell.merge(shape.cell(row_index, len(headers) - 1))
                cell.text = values["section"]
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.color("E1F1ED")
                continue
            for column, value in enumerate(values):
                cell = shape.cell(row_index, column)
                cell.text = str(value)
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.color("F4F9F8" if row_index % 2 == 0 else "FFFFFF")
        for row in shape.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size, paragraph.font.name = self.pt(font_size), "Arial"
                    paragraph.font.color.rgb = self.color(INK)
        for cell in shape.rows[0].cells:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.color.rgb = self.color("FFFFFF")
        for row_index in sections:
            for paragraph in shape.cell(row_index, 0).text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = self.pt(font_size + 1)
                paragraph.font.color.rgb = self.color(TEAL)

    def paginated_table(self, title, headers, rows, widths, subtitle=None, empty_note=None, rows_per_slide=15, font_size=8, sections=None):
        """Carry a full register across as many slides as it needs.

        ``sections`` names the specification category of each row; the category heading is
        written above its rows and repeated whenever a category continues on a new slide.
        """
        offset = 1.45 if subtitle is None else 1.72
        if not rows:
            slide = self.new_slide(title)
            if subtitle:
                self.add_text(slide, subtitle, .42, 1.24, 12.2, .26, 12, GREY, wrap=True)
            self.table(slide, headers, [], widths, y=offset, font_size=font_size, empty_note=empty_note or "No lines fall in this register for the selected reporting date.")
            return
        if sections:
            rows_per_slide = min(rows_per_slide, 13)
        for start in range(0, len(rows), rows_per_slide):
            end = min(start + rows_per_slide, len(rows))
            suffix = f" ({start + 1}–{end} of {len(rows)})" if len(rows) > rows_per_slide else ""
            slide = self.new_slide(f"{title}{suffix}")
            if subtitle:
                self.add_text(slide, subtitle, .42, 1.24, 12.2, .26, 12, GREY, wrap=True)
            page = rows[start:end]
            if sections:
                page, current = [], None
                carried = sections[start - 1] if start else None
                for row, section in zip(rows[start:end], sections[start:end]):
                    if section != current:
                        page.append({"section": f"{section} (continued)" if current is None and section == carried else section})
                        current = section
                    page.append(row)
            self.table(slide, headers, page, widths, y=offset, font_size=font_size)

    def save(self, filename: str) -> tuple[BytesIO, str]:
        output = BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output, filename


def _scope_labels(data: dict[str, Any], centre_ids: set[int] | None) -> tuple[str, str, str]:
    """How a deck names its own scope: heading, one-line description, filename stem."""
    if centre_ids is None:
        return "all ONGC", "all ONGC", "ONGC"
    names = [centre["name"] for centre in data.get("scope_centres") or []]
    zones = sorted({centre["zone"] for centre in data.get("scope_centres") or [] if centre["zone"]})
    if not names:
        return "selected assets", "the selected assets", "Selected assets"
    if len(names) == 1:
        return names[0], names[0], names[0]
    listed = ", ".join(names[:4]) + (f" and {len(names) - 4} more" if len(names) > 4 else "")
    heading = zones[0] if len(zones) == 1 else f"{len(names)} assets"
    return heading, f"{len(names)} assets — {listed}", heading


def build_management_review_presentation(static_folder: str, reporting_date: date | None = None, compare_date: date | None = None, centre_ids: set[int] | None = None) -> tuple[BytesIO, str]:
    """Build the management-review deck for one published reporting date.

    ``centre_ids`` builds it for a chosen set of assets — a zone, a group of
    assets, or one — instead of the whole of ONGC. Every figure in the deck then
    counts only those assets, and the deck says so on its cover.
    """
    from app.core.services.inventory_monitoring import management_review_data

    data = management_review_data(reporting_date, compare_date, centre_ids)
    if data["reporting_date"] is None or data["kpis"] is None:
        raise ValueError(
            "No published Group 09 and Group 10 stock was found for the selected assets."
            if centre_ids is not None else
            "Publish Group 09 and Group 10 workbooks for a reporting date before downloading the presentation."
        )
    kpis, selected, previous, thresholds = data["kpis"], data["reporting_date"], data["previous_date"], data["thresholds"]
    period_label = selected.strftime("%d %B %Y")
    scope_heading, scope_note, scope_stem = _scope_labels(data, centre_ids)
    deck = _Deck(static_folder, f"Source: Published Group 09 and Group 10 inventory snapshots as on {period_label} · Scope: {scope_note}")

    def delta_note(current, prior, positive_is_bad=True):
        if prior is None:
            return "First published period · no prior comparison", GREY
        change = (current or Decimal("0")) - (prior or Decimal("0"))
        if change == 0:
            return "No change against the comparison period", GREY
        direction = "Increase" if change > 0 else "Reduction"
        tone = (RED if change > 0 else TEAL) if positive_is_bad else (TEAL if change > 0 else RED)
        return f"{direction} of {_crore(abs(change))} against {previous:%d %b %Y}", tone

    def register_rows(register):
        return [
            [row["code"], row.get("uom") or "—", _concise(row["description"], 40), row["group"], _concise(row["centre"], 30), _concise(row["zone"], 20),
             _quantity(row["qty"], None), _crore(row["value"]), _months(row["months"])]
            for row in register["rows"]
        ]

    def register_subtitle(register):
        note = f"{register['description']}  ·  {register['count']:,} stock lines  ·  {_crore(register['value'])} of inventory value"
        if register["omitted"]:
            note += f"  ·  {register['omitted']:,} further lines beyond the first {data['group_limit']} of each category are in the application register"
        return note

    deck.cover(
        f"INVENTORY MONITORING · GROUPS 09 AND 10 · {scope_heading.upper()}", "Inventory Management Review", f"Position as on {period_label}",
        f"{_crore(kpis['total_value'])} of monitored chemicals — Groups 09 and 10 counted together — across {kpis['centre_count']:,} work centres and {kpis['material_count']:,} materials, {scope_note}",
    )

    slide = deck.new_slide("Executive summary")
    risk_note, risk_tone = delta_note(kpis["value_at_risk"], kpis["prev_value_at_risk"])
    stockout_note, stockout_tone = delta_note(kpis["stockout_value"], kpis["prev_stockout_value"])
    total_note, _tone = delta_note(kpis["total_value"], kpis["prev_total_value"])
    cards = [
        (_crore(kpis["total_value"]), "Total monitored inventory", INK, total_note),
        (_crore(kpis["value_at_risk"]), "Working capital at risk", risk_tone if kpis["prev_value_at_risk"] is not None else AMBER, f"{kpis['at_risk_share']}% in slow-moving or excess bands"),
        (_crore(kpis["stockout_value"]), "Stock-out exposure", RED, f"{kpis['stockout_share']}% at or below low-stock coverage"),
        (f"{len(data['high_value_materials']):,}", "Materials above ₹ 1 Cr", TEAL, f"{_crore(data['high_value_total'])} held in these materials"),
        (f"{kpis['top5_share']}%", "Held by top five work centres", INK, f"{kpis['centre_count']:,} work centres reporting stock"),
        (f"{kpis['material_count']:,}", "Materials reporting stock", TEAL, f"Across {kpis['record_count']:,} stock lines in this period"),
    ]
    for index, (value, label, tone, note) in enumerate(cards):
        deck.metric(slide, .6 + (index % 3) * 4.2, 1.62 + (index // 3) * 2.15, value, label, tone, note)
    if previous:
        deck.add_text(slide, risk_note, .6, 6.15, 12.1, .3, 13, risk_tone, wrap=True)
        deck.add_text(slide, stockout_note, .6, 6.48, 12.1, .3, 13, stockout_tone, wrap=True)
    else:
        deck.add_text(slide, "This is the first published reporting period; movement analysis begins from the next published period.", .6, 6.2, 12.1, .3, 13, GREY, wrap=True)

    slide = deck.new_slide("Stock health by coverage band")
    deck.add_text(slide, f"Thresholds in force: critical at or below {_months(thresholds['critical_low_stock_months'])} months, low at or below {_months(thresholds['low_stock_months'])} months, slow-moving from {_months(thresholds['slow_moving_months'])} months, excess from {_months(thresholds['excess_stock_months'])} months.", .42, 1.24, 12.2, .3, 12, GREY, wrap=True)
    deck.table(slide, ["Coverage band", "Stock lines", "Inventory value", "Share of value"], [
        [band["label"], f"{band['count']:,}", _crore(band["value"]), f"{band['share']}%"] for band in data["health_mix"]
    ], [4.35, 2.5, 3.1, 2.5], y=1.75, font_size=11, empty_note="No stock lines were classified for this reporting date.")

    slide = deck.new_slide("Inventory value by zone")
    deck.table(slide, ["Zone", "Inventory value", "Share of portfolio", "Comparison period", "Movement"], [
        [zone["zone"], _crore(zone["value"]), f"{zone['share']}%",
         _crore(zone["prev"]) if zone["prev"] is not None else "—",
         (_crore((zone["value"] or Decimal('0')) - (zone["prev"] or Decimal('0'))) if zone["prev"] is not None else "First period")]
        for zone in data["zones"]
    ], [3.5, 2.6, 2.35, 2.1, 1.9], font_size=10, empty_note="No zone-level inventory was reported for this date.")

    deck.paginated_table(
        "Work centres by inventory value",
        ["#", "Work centre", "Zone", "Inventory value", "Share of portfolio"],
        [[index, centre["name"], centre["zone"], _crore(centre["value"]), f"{centre['share']}%"] for index, centre in enumerate(data["centres_ranked"], 1)],
        [.5, 4.6, 3.2, 2.2, 1.95],
        subtitle="Every work centre holding monitored Group 09 or Group 10 stock, ranked by value.",
        empty_note="No work centre reported monitored stock for this date.",
    )

    comparison = data.get("comparison")
    slide = deck.new_slide("Movement against the comparison period")
    if comparison:
        deck.add_text(slide, f"Like-for-like movement over {comparison['gap_days']:,} days: {comparison['common_centres']:,} work centres reported stock in both {previous:%d %b %Y} and {selected:%d %b %Y}.", .42, 1.24, 12.2, .28, 12, GREY, wrap=True)
        deck.add_text(slide, "Largest build-ups", .8, 1.8, 5.2, .3, 17, RED, True)
        for index, mover in enumerate(data["movers"]["up"] or [{"name": "No work centre increased its holding.", "delta": None}]):
            deck.add_text(slide, mover["name"], .8, 2.25 + index * .45, 4.4, .28, 13, INK, wrap=True)
            deck.add_text(slide, _crore(mover["delta"]) if mover["delta"] is not None else "", 5.3, 2.25 + index * .45, 1.5, .28, 13, RED, True)
        deck.add_text(slide, "Largest draw-downs", 7.1, 1.8, 5.2, .3, 17, TEAL, True)
        for index, mover in enumerate(data["movers"]["down"] or [{"name": "No work centre reduced its holding.", "delta": None}]):
            deck.add_text(slide, mover["name"], 7.1, 2.25 + index * .45, 4.4, .28, 13, INK, wrap=True)
            deck.add_text(slide, _crore(-mover["delta"]) if mover["delta"] is not None else "", 11.6, 2.25 + index * .45, 1.5, .28, 13, TEAL, True)
        net = comparison["like_for_like_delta"]
        change = f" ({comparison['like_for_like_change']:+.1f}%)" if comparison["like_for_like_change"] is not None else ""
        for index, item in enumerate([
            (_crore(comparison["like_for_like"]), f"Like-for-like, {selected:%d %b %Y}", INK, f"{comparison['common_centres']:,} work centres in both periods"),
            (_crore(comparison["like_for_like_prev"]), f"Like-for-like, {previous:%d %b %Y}", GREY, f"Comparison period, {comparison['gap_days']:,} days earlier"),
            (f"{_crore(net)}{change}", "Net like-for-like movement", RED if net > 0 else TEAL, "Excludes work centres reporting in only one period"),
        ]):
            deck.metric(slide, .8 + index * 4.2, 4.55, item[0], item[1], item[2], item[3])
        scope = []
        if comparison["entrant_count"]:
            scope.append(f"{comparison['entrant_count']:,} work centres report stock only this period ({_crore(comparison['entrant_value'])})")
        if comparison["exit_count"]:
            scope.append(f"{comparison['exit_count']:,} reported only in the comparison period ({_crore(comparison['exit_value'])})")
        deck.add_text(slide, ("Scope difference: " + "; ".join(scope) + ".") if scope else "Both periods cover exactly the same work centres.", .8, 6.45, 11.7, .3, 12, AMBER if scope else TEAL, wrap=True)
    else:
        deck.add_text(slide, "This is the first published reporting period. Movement analysis is available once a second period is published.", .8, 2.1, 11.2, .4, 17, GREY, wrap=True)

    deck.paginated_table(
        "Materials above ₹ 1 Cr of inventory value",
        ["Material", "UoM", "Description", "Grp", "Inventory value", "Share", "Work centres", "Coverage range"],
        [[item["code"], item.get("uom") or "—", _concise(item["description"], 54), item["group"], _crore(item["value"]), f"{item['share']}%", f"{item['centres']:,}",
          "—" if item["months_low"] is None else f"{_months(item['months_low'])} – {_months(item['months_high'])} months"]
         for item in data["high_value_materials"]],
        [1.4, .55, 3.95, .6, 1.6, .9, 1.1, 2.35],
        sections=[item["section"] for item in data["high_value_materials"]],
        subtitle=f"All-ONGC holdings for every material whose total value crosses ₹ 1 Cr, in corporate specification order · {_crore(data['high_value_total'])} across {len(data['high_value_materials']):,} materials.",
        empty_note="No single material crosses ₹ 1 Cr of inventory value at this reporting date.",
    )

    for register in data["coverage_registers"]:
        deck.paginated_table(
            register["label"], REGISTER_HEADERS, register_rows(register), REGISTER_WIDTHS,
            subtitle=register_subtitle(register), sections=[row["section"] for row in register["rows"]],
        )
    for register in data["supporting_registers"]:
        deck.paginated_table(
            register["label"], SUPPORTING_HEADERS,
            [[row["code"], row.get("uom") or "—", _concise(row["description"], 48), row["group"] or "—", _concise(row["centre"], 32), _concise(row["zone"], 22), _crore(row["value"])] for row in register["rows"]],
            SUPPORTING_WIDTHS,
            sections=[row["section"] for row in register["rows"]],
            subtitle=register_subtitle(register).replace("stock lines", "material lines"),
            empty_note=f"The source workbook for {period_label} reported no items in this register.",
        )

    slide = deck.new_slide("Decisions sought from the review")
    non_moving = next((item for item in data["supporting_registers"] if item["key"] == "non_moving"), {"count": 0, "value": Decimal("0")})
    excess = next((item for item in data["coverage_registers"] if item["key"] == "excess_stock"), {"count": 0, "value": Decimal("0")})
    critical = next((item for item in data["coverage_registers"] if item["key"] == "critical_low_stock"), {"count": 0, "value": Decimal("0")})
    open_supply = next((item for item in data["coverage_registers"] if item["key"] == "open_supply_with_high_stock"), {"count": 0, "value": Decimal("0")})
    for index, item in enumerate([
        (_crore(excess["value"]), "Excess stock to be redeployed", PURPLE, f"{excess['count']:,} stock lines at or above the excess threshold"),
        (_crore(non_moving["value"]), "Non-moving stock for disposal review", AMBER, f"{non_moving['count']:,} material lines in the source register"),
        (_crore(critical["value"]), "Critical low stock to be replenished", RED, f"{critical['count']:,} stock lines at or below the critical threshold"),
    ]):
        deck.metric(slide, .8 + index * 4.2, 1.7, item[0], item[1], item[2], item[3])
    deck.add_text(slide, "Points for direction", .8, 3.7, 6.0, .3, 18, INK, True)
    for index, question in enumerate([
        f"Which work centres will absorb the {_crore(excess['value'])} of excess stock, and by when?",
        f"Which of the {non_moving['count']:,} non-moving lines are approved for transfer, use-up or disposal?",
        f"Which of the {open_supply['count']:,} open PO / PR lines against high stock should be held, reduced or cancelled ({_crore(open_supply['value'])} of stock affected)?",
        f"What replenishment action is committed for the {critical['count']:,} critical low-stock lines?",
        "Which work-centre / material pairs need mapping correction before the next reporting date?",
    ]):
        deck.add_text(slide, f"{index + 1}", .8, 4.15 + index * .48, .25, .25, 15, TEAL, True)
        deck.add_text(slide, question, 1.2, 4.15 + index * .48, 11.4, .3, 13, INK, wrap=True)

    return deck.save(f"{scope_stem} Inventory Management Review {selected:%d %b %Y}.pptx")


CENTRE_BAND_SLIDES = (
    ("critical_low_stock", "Critical low stock", "At or below the critical threshold — replenishment action is required now."),
    ("low_stock", "Low stock", "Below the low-stock threshold — plan replenishment."),
    ("slow_moving_stock", "Slow-moving stock", "At or above the slow-moving threshold — review consumption and transfers."),
    ("excess_stock", "Excess stock", "At or above the excess threshold — stop-buy or redeploy."),
)


def build_work_centre_review_presentation(static_folder: str, work_center_id: int, unit: str | None = None, compare_date: date | None = None) -> tuple[BytesIO, str]:
    """Build the review deck for a single work centre, on the same registers the page shows."""
    from app.core.services.inventory_monitoring import work_center_review_data

    data = work_center_review_data(work_center_id, unit, compare_date)
    centre, kpis, thresholds = data["centre"], data["kpis"], data["thresholds"]
    selected, previous, comparison = data["reporting_date"], data["previous_date"], data["comparison"]
    if selected is None or not kpis["line_count"]:
        raise ValueError("This work centre has no mapped stock lines yet, so there is nothing to present.")
    unit_label = data["selected_unit"] or "Combined asset status"
    period_label = selected.strftime("%d %B %Y")
    deck = _Deck(static_folder, f"Source: {centre.name} · {unit_label} · mapped Group 09 and Group 10 stock as on {period_label}")

    quantities: dict[str, tuple[Any, str | None]] = {}
    for rows in data["groups"].values():
        for record in rows:
            held_qty, uom = quantities.get(record.material_code, (None, None))
            quantities[record.material_code] = (
                (record.stock_qty or Decimal("0")) + (held_qty or Decimal("0")) if record.stock_qty is not None or held_qty is not None else None,
                uom or record.uom,
            )

    def band_rows(key):
        """Rows and their specification category, in DFC / 01, DFC / 02 … order."""
        rows, sections = [], []
        for group in data["spec_groups"][key]:
            for record in group["rows"]:
                rows.append([
                    record.material_code,
                    (record.material.uom if record.material and record.material.uom else record.uom) or "—",
                    _concise(record.material_description, 60), record.material_group,
                    _quantity(record.stock_qty, None), _crore(record.inventory_value_inr), _months(record.stock_months),
                ])
                sections.append(group["label"])
        return rows, sections

    deck.cover(
        f"WORK CENTRE REVIEW · {(centre.zone or 'Unassigned zone').upper()}", centre.name, f"Position as on {period_label}",
        f"{_crore(kpis['total_value'])} of mapped chemicals — Groups 09 and 10 counted together — across {kpis['material_count']:,} materials · {unit_label} · {kpis['portfolio_share']}% of the all-ONGC monitored portfolio",
    )

    slide = deck.new_slide(f"{centre.name} | Position summary")
    cards = [
        (_crore(kpis["total_value"]), "Inventory held at this centre", INK, f"{kpis['line_count']:,} stock lines · {kpis['material_count']:,} materials"),
        (_crore(kpis["value_at_risk"]), "Working capital at risk", AMBER, f"{kpis['at_risk_share']}% in slow-moving or excess bands"),
        (_crore(kpis["stockout_value"]), "Stock-out exposure", RED, f"{kpis['stockout_share']}% at or below low-stock coverage"),
        (f"{kpis['portfolio_share']}%", "Share of all-ONGC portfolio", TEAL, f"Against {_crore(kpis['portfolio_value'])} monitored in the same workbooks"),
        (f"{len(data['groups']['critical_low_stock']):,}", "Critical low-stock lines", RED, f"{len(data['groups']['excess_stock']):,} lines in the excess band"),
        (f"{kpis['source_case_count']:,}", "Source-reported cases", AMBER, f"{_crore(kpis['source_case_value'])} in non-moving, aged, surplus or transit registers"),
    ]
    for index, (value, label, tone, note) in enumerate(cards):
        deck.metric(slide, .6 + (index % 3) * 4.2, 1.62 + (index // 3) * 2.15, value, label, tone, note)
    dates = {value for value in data["as_on_by_group"].values() if value}
    mixed = "; ".join(f"Group {group} as on {value:%d %b %Y}" for group, value in sorted(data["as_on_by_group"].items()) if value)
    deck.add_text(
        slide,
        f"Mixed source dates in this view — {mixed}." if len(dates) > 1 else "Groups 09 and 10 are counted together throughout this deck; every stock line reported at this work centre is included.",
        .6, 6.2, 12.1, .3, 12, AMBER if len(dates) > 1 else GREY, wrap=True,
    )

    slide = deck.new_slide("Stock health by coverage band")
    deck.add_text(slide, f"Thresholds in force: critical at or below {_months(thresholds['critical_low_stock_months'])} months, low at or below {_months(thresholds['low_stock_months'])} months, slow-moving from {_months(thresholds['slow_moving_months'])} months, excess from {_months(thresholds['excess_stock_months'])} months.", .42, 1.24, 12.2, .3, 12, GREY, wrap=True)
    deck.table(slide, ["Coverage band", "Stock lines", "Inventory value", "Share of centre value", "Comparison period"], [
        [band["label"], f"{band['count']:,}", _crore(band["value"]), f"{band['share']}%",
         _crore(band["prev"]) if band["prev"] is not None else "—"]
        for band in data["health_mix"]
    ], [3.6, 2.2, 2.6, 2.35, 1.7], y=1.9, font_size=12, empty_note="No mapped stock lines were classified for this work centre.")

    slide = deck.new_slide("Movement against the comparison period")
    if comparison:
        deck.add_text(slide, f"Like-for-like movement over {comparison['gap_days']:,} days: {comparison['common_materials']:,} materials were held at this centre on both {previous:%d %b %Y} and {selected:%d %b %Y}.", .42, 1.24, 12.2, .28, 12, GREY, wrap=True)
        deck.add_text(slide, "Largest build-ups", .8, 1.8, 5.2, .3, 17, RED, True)
        for index, mover in enumerate(data["movers"]["up"] or [{"code": "", "description": "No material increased its holding.", "delta": None}]):
            deck.add_text(slide, f"{mover['code']} {_concise(mover['description'], 34)}".strip(), .8, 2.25 + index * .45, 4.4, .28, 12, INK, wrap=True)
            deck.add_text(slide, _crore(mover["delta"]) if mover["delta"] is not None else "", 5.3, 2.25 + index * .45, 1.5, .28, 12, RED, True)
        deck.add_text(slide, "Largest draw-downs", 7.1, 1.8, 5.2, .3, 17, TEAL, True)
        for index, mover in enumerate(data["movers"]["down"] or [{"code": "", "description": "No material reduced its holding.", "delta": None}]):
            deck.add_text(slide, f"{mover['code']} {_concise(mover['description'], 34)}".strip(), 7.1, 2.25 + index * .45, 4.4, .28, 12, INK, wrap=True)
            deck.add_text(slide, _crore(-mover["delta"]) if mover["delta"] is not None else "", 11.6, 2.25 + index * .45, 1.5, .28, 12, TEAL, True)
        net = comparison["like_for_like_delta"]
        change = f" ({comparison['like_for_like_change']:+.1f}%)" if comparison["like_for_like_change"] is not None else ""
        for index, item in enumerate([
            (_crore(comparison["like_for_like"]), f"Like-for-like, {selected:%d %b %Y}", INK, f"{comparison['common_materials']:,} materials in both periods"),
            (_crore(comparison["like_for_like_prev"]), f"Like-for-like, {previous:%d %b %Y}", GREY, f"Comparison period, {comparison['gap_days']:,} days earlier"),
            (f"{_crore(net)}{change}", "Net like-for-like movement", RED if net > 0 else TEAL, "Excludes materials held in only one period"),
        ]):
            deck.metric(slide, .8 + index * 4.2, 4.55, item[0], item[1], item[2], item[3])
        scope = []
        if comparison["entrant_count"]:
            scope.append(f"{comparison['entrant_count']:,} materials appear only now ({_crore(comparison['entrant_value'])})")
        if comparison["exit_count"]:
            scope.append(f"{comparison['exit_count']:,} were held only on {previous:%d %b %Y} ({_crore(comparison['exit_value'])})")
        deck.add_text(slide, ("Scope difference: " + "; ".join(scope) + ".") if scope else "Both periods cover exactly the same materials.", .8, 6.45, 11.7, .3, 12, AMBER if scope else TEAL, wrap=True)
    else:
        deck.add_text(slide, "Only one reporting date is available for this work centre. Movement analysis begins from the next import.", .8, 2.1, 11.2, .4, 17, GREY, wrap=True)

    deck.paginated_table(
        "Materials held", ["Material", "UoM", "Description", "Grp", "Stock qty", "Inventory value", "Share of centre"],
        [[item["code"], quantities.get(item["code"], (None, None))[1] or "—", _concise(item["description"], 54), item["group"],
          _quantity(quantities.get(item["code"], (None, None))[0], None), _crore(item["value"]), f"{item['share']}%"]
         for item in data["top_materials"]],
        [1.5, .55, 4.05, .6, 1.85, 1.95, 1.95],
        sections=[item["section"] for item in data["top_materials"]],
        subtitle=f"Every mapped material held at {centre.name}, in corporate specification order.",
        empty_note="No mapped material is held at this work centre.",
    )

    for key, label, description in CENTRE_BAND_SLIDES:
        rows, sections = band_rows(key)
        deck.paginated_table(
            label, CENTRE_REGISTER_HEADERS, rows, CENTRE_REGISTER_WIDTHS, sections=sections,
            subtitle=f"{description}  ·  {len(rows):,} materials  ·  {_crore(sum((record.inventory_value_inr or Decimal('0') for record in data['groups'][key]), Decimal('0')))}",
            empty_note=f"No material at {centre.name} falls in the {label.lower()} band.",
        )

    source_rows, source_sections = [], []
    for group in data["source_spec_groups"]:
        for item in group["rows"]:
            source_rows.append([
                item.exception_type.replace("_", " ").title(),
                item.material.material_code if item.material else "—",
                (item.material.uom if item.material else None) or "—",
                _concise(item.material.description if item.material else None, 46),
                _concise(item.details, 62),
            ])
            source_sections.append(group["label"])
    deck.paginated_table(
        "Non-moving, aged, surplus and transit cases",
        ["Condition", "Material", "UoM", "Chemical / material name", "Source evidence"],
        source_rows, [2.2, 1.5, .55, 3.8, 4.4], sections=source_sections,
        subtitle="Reported directly by the source inventory workbook for this work centre, not inferred from stock months.",
        empty_note=f"The source workbook reports no non-moving, aged, surplus or transit case at {centre.name}.",
    )

    slide = deck.new_slide(f"Decisions sought for {centre.name}")
    critical, excess = data["groups"]["critical_low_stock"], data["groups"]["excess_stock"]
    non_moving = next((item for item in data["source_summary"] if item["key"] == "non_moving"), {"count": 0, "value": Decimal("0")})
    for index, item in enumerate([
        (_crore(sum((record.inventory_value_inr or Decimal("0") for record in excess), Decimal("0"))), "Excess stock to be redeployed", PURPLE, f"{len(excess):,} materials at or above the excess threshold"),
        (_crore(non_moving["value"]), "Non-moving stock for disposal review", AMBER, f"{non_moving['count']:,} materials in the source register"),
        (_crore(sum((record.inventory_value_inr or Decimal("0") for record in critical), Decimal("0"))), "Critical low stock to be replenished", RED, f"{len(critical):,} materials at or below the critical threshold"),
    ]):
        deck.metric(slide, .8 + index * 4.2, 1.7, item[0], item[1], item[2], item[3])
    deck.add_text(slide, "Points for direction", .8, 3.7, 6.0, .3, 18, INK, True)
    for index, question in enumerate([
        f"Which of the {len(critical):,} critical low-stock materials at {centre.name} have a committed replenishment date?",
        f"Can the {len(excess):,} excess-stock materials be transferred to a centre that consumes them?",
        f"Which of the {non_moving['count']:,} non-moving materials are approved for use-up, transfer or disposal?",
        "Are the coverage figures consistent with the drilling programme planned for this asset?",
        "Which mapped materials are no longer required at this work centre and should leave the mapping?",
    ]):
        deck.add_text(slide, f"{index + 1}", .8, 4.15 + index * .48, .25, .25, 15, TEAL, True)
        deck.add_text(slide, question, 1.2, 4.15 + index * .48, 11.4, .3, 13, INK, wrap=True)

    suffix = f" {data['selected_unit']}" if data["selected_unit"] else ""
    return deck.save(f"{centre.name}{suffix} Inventory Review {selected:%d %b %Y}.pptx")
