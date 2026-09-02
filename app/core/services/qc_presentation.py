"""Lazy, on-demand PowerPoint exports for QC laboratory dashboards."""
from __future__ import annotations

from datetime import date, timedelta
from io import BytesIO
from pathlib import Path

from app.models.quality_control.qc_sample import QCSample
from app.models.quality_control.qc_testing_standard import QCTestingStandard


# Deck chrome shared by every QC export: one navy rule at the top, a thin border
# above and below the body, the ONGC and Corporate Chemistry marks, and a footer
# carrying the data source.
# Kept module-level so the lab decks cannot drift apart visually.
class _DeckChrome:
    """Slide furniture for a QC deck, bound to one Presentation."""

    NAVY, BLUE, RED, GREEN, GREY = "071D42", "1976D2", "C53B3B", "18794E", "526173"
    BORDER = "D7E2EE"

    def __init__(self, prs, static_folder: str, source_line: str):
        from pptx.util import Inches, Pt
        self.prs = prs
        self.blank = prs.slide_layouts[6]
        self.source_line = source_line
        self.corporate_chemistry_logo = Path(static_folder) / "images" / "ongc-corporate-chemistry-logo.png"
        self.ongc_logo = Path(static_folder) / "images" / "ongc-official-logo.png"
        self._Inches, self._Pt = Inches, Pt

    def color(self, value):
        from pptx.dml.color import RGBColor
        return RGBColor.from_string(value)

    def add_text(self, slide, value, x, y, w, h, size=18, tone=None, bold=False):
        from pptx.util import Inches, Pt
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        p = shape.text_frame.paragraphs[0]
        p.text = str(value)
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = self.color(tone or self.NAVY)
        p.font.name = "Arial"
        return shape

    def add_wrapped_text(self, slide, value, x, y, w, h, size=18, tone=None, bold=False):
        shape = self.add_text(slide, value, x, y, w, h, size, tone, bold)
        shape.text_frame.word_wrap = True
        return shape

    def rectangle(self, slide, x, y, w, h, fill, line=None):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color(fill)
        shape.line.color.rgb = self.color(line or fill)
        return shape

    def canvas_background(self, slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.color("FFFFFF")
        self.rectangle(slide, 0, 0, 13.333, .08, self.NAVY)
        self.rectangle(slide, .42, 1.26, 12.45, .015, self.BORDER)
        self.rectangle(slide, .42, 6.93, 12.45, .015, self.BORDER)

    def add_header_branding(self, slide):
        """Place both approved marks in the deck header without crowding it."""
        if self.ongc_logo.exists():
            slide.shapes.add_picture(str(self.ongc_logo), self._Inches(11.15), self._Inches(.21), width=self._Inches(.9), height=self._Inches(.51))
        if self.corporate_chemistry_logo.exists():
            slide.shapes.add_picture(str(self.corporate_chemistry_logo), self._Inches(12.24), self._Inches(.16), width=self._Inches(.55), height=self._Inches(.55))

    def add_cover_branding(self, slide, background="FFFFFF"):
        """Use the same ONGC + Corporate Chemistry pairing on title slides."""
        self.rectangle(slide, 10.0, .46, 2.32, 1.2, background, self.BORDER)
        if self.ongc_logo.exists():
            slide.shapes.add_picture(str(self.ongc_logo), self._Inches(10.16), self._Inches(.68), width=self._Inches(.9), height=self._Inches(.51))
        if self.corporate_chemistry_logo.exists():
            slide.shapes.add_picture(str(self.corporate_chemistry_logo), self._Inches(11.16), self._Inches(.53), width=self._Inches(1.05), height=self._Inches(1.05))

    def header(self, slide, title, page):
        from pptx.util import Inches
        self.canvas_background(slide)
        self.add_text(slide, "ONGC CORPORATE CHEMISTRY \u00b7 QC LABORATORY MONITORING", .42, .27, 7.4, .24, 11, self.BLUE, True)
        self.add_text(slide, title, .42, .7, 11.5, .46, 27, self.NAVY, True)
        self.add_header_branding(slide)
        self.add_text(slide, self.source_line, .42, 7.08, 8.6, .16, 8, self.GREY)
        self.add_text(slide, f"{page:02d}", 12.5, 7.08, .25, .16, 8, self.GREY)

    def metric(self, slide, x, y, value, label, tone=None):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches
        tone = tone or self.NAVY
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - .18), Inches(y - .18), Inches(3.45), Inches(1.35))
        card.fill.solid()
        card.fill.fore_color.rgb = self.color(
            "EAF4FF" if tone == self.BLUE else "FCEBEC" if tone == self.RED
            else "EAF7F0" if tone == self.GREEN else "F2F4F7"
        )
        card.line.color.rgb = self.color(self.BORDER)
        self.add_text(slide, value, x, y, 2.7, .4, 27, tone, True)
        self.add_text(slide, label, x, y + .52, 3.0, .25, 13, self.NAVY, True)

    def new_slide(self, title, page):
        slide = self.prs.slides.add_slide(self.blank)
        self.header(slide, title, page)
        return slide


def _paginated_rows(rows, page_size: int):
    """Split an operational register without ever dropping the final rows."""
    if page_size < 1:
        raise ValueError("A presentation page must contain at least one row.")
    return [rows[index:index + page_size] for index in range(0, len(rows), page_size)] or [[]]


def _sap_presentation_action_groups(data):
    """Order SAP-open samples by laboratory, then corporate-specification group.

    The all-laboratories deck is read as a hand-off pack. A mixed, globally
    overdue-sorted register makes it difficult for an RGL to find its own
    Drilling Fluid, Production Chemical and other queues. Keep the approved
    SAP laboratory order first, followed by the Corporate Specification
    category order; overdue items still lead within each section.
    """
    from app.core.services.csc_utils import SPEC_SUBSET_ORDER
    from app.core.services.sap_quality_control import CORPORATE_SPECIFICATION_UNMATCHED_KEY

    laboratory_rank = {
        laboratory["code"]: index
        for index, laboratory in enumerate(data["scope_laboratories"])
    }
    subgroup_rank = {key: index for index, key in enumerate(SPEC_SUBSET_ORDER)}
    grouped: dict[tuple[str, str], list] = {}
    labels: dict[tuple[str, str], str] = {}
    laboratories: dict[str, dict] = {}

    for entry in data["action_entries"]:
        laboratory = entry["laboratory"]
        subgroup_key = entry["subgroup_key"] or CORPORATE_SPECIFICATION_UNMATCHED_KEY
        key = (laboratory["code"], subgroup_key)
        grouped.setdefault(key, []).append(entry)
        labels[key] = entry["subgroup_label"]
        laboratories[laboratory["code"]] = laboratory

    groups = []
    for key in sorted(
        grouped,
        key=lambda value: (
            laboratory_rank.get(value[0], len(laboratory_rank)),
            99 if value[1] == CORPORATE_SPECIFICATION_UNMATCHED_KEY else subgroup_rank.get(value[1], 90),
            value[1],
        ),
    ):
        entries = sorted(
            grouped[key],
            key=lambda entry: (
                not entry["stt_overdue"],
                entry["stt_due_date"] or date.max,
                entry["specification_no"] or "",
                (entry["record"].material_description or "").casefold(),
                entry["record"].notification_no or "",
                entry["record"].id,
            ),
        )
        groups.append({
            "laboratory": laboratories[key[0]],
            "subgroup_key": key[1],
            "subgroup_label": labels[key],
            "entries": entries,
        })
    return groups


def build_lab_performance_presentation(lab_code: str, static_folder: str) -> tuple[BytesIO, str]:
    """Create a lab-specific review deck without loading PPTX libraries at startup."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
    from app.core.services.quality_control import CLOSED_SAMPLE_REVIEW_STT_DAYS, _normalized_chemical, latest_dashboard_data

    data = latest_dashboard_data(lab_code)
    batch = data["batch"]
    if batch is None:
        raise ValueError("Import a local status workbook before downloading a presentation.")
    standards = {item.normalized_name: item.standard_days for item in QCTestingStandard.query.all()}
    month_start = batch.week_end.replace(day=1)
    next_month = (month_start.replace(day=28) + timedelta(days=4)).replace(day=1)
    completed = QCSample.query.filter(QCSample.lab_code == lab_code, QCSample.report_issue_date >= month_start, QCSample.report_issue_date < next_month).all()
    completed = [s for s in completed if s.result_status in {"pass", "fail", "report_issued"} and s.turnaround_days is not None]
    def stt(sample): return standards.get(_normalized_chemical(sample.chemical_name)) or CLOSED_SAMPLE_REVIEW_STT_DAYS
    late, within = [s for s in completed if s.turnaround_days > stt(s)], [s for s in completed if s.turnaround_days <= stt(s)]
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    chrome = _DeckChrome(prs, static_folder, f"Source: {data['laboratory']['name']} QC data \u00b7 {data['month_label']}")
    blank = chrome.blank
    navy, blue, red, green, grey = chrome.NAVY, chrome.BLUE, chrome.RED, chrome.GREEN, chrome.GREY
    border = chrome.BORDER
    color, add_text, add_wrapped_text = chrome.color, chrome.add_text, chrome.add_wrapped_text
    rectangle, header, metric = chrome.rectangle, chrome.header, chrome.metric
    canvas_background = chrome.canvas_background
    def overview_text(value, limit=74):
        """Keep narrative spreadsheet remarks readable in overview slides."""
        text = " ".join(str(value or "Not recorded").split())
        if len(text) <= limit:
            return text
        return f"{text[:limit + 1].rsplit(' ', 1)[0].rstrip('.,;:')}…"
    def table_slide(title, rows, page, reason=False):
        slide = prs.slides.add_slide(blank); header(slide, title, page)
        if not rows:
            add_text(slide, "No completed samples fall in this category for the selected calendar month.", .8, 2.0, 10.5, .35, 18, green, True)
            return
        if len(rows) == 1:
            sample = rows[0]
            add_text(slide, "Sample detail", .8, 1.5, 3.5, .3, 18, navy, True)
            rectangle(slide, .8, 1.95, 11.6, 2.25, "FFFFFF", border)
            rectangle(slide, .8, 1.95, 11.6, .06, red if reason else green)
            add_text(slide, sample.chemical_name, 1.1, 2.3, 5.2, .35, 24, navy, True)
            add_text(slide, f"Notification: {sample.notification_no or sample.po_number or '—'}", 1.1, 2.78, 3.8, .25, 14, grey)
            add_text(slide, f"Received: {sample.sample_receipt_date:%d %b %Y} · Reported: {sample.report_issue_date:%d %b %Y}", 5.0, 2.78, 4.8, .25, 14, grey)
            add_text(slide, f"Actual TAT: {sample.turnaround_days} days · STT: {stt(sample)} days · Variance: {sample.turnaround_days-stt(sample):+d} days", 1.1, 3.2, 6.8, .25, 15, red if reason else green, True)
            add_text(slide, f"Delay reason: {sample.delay_reason or 'Not recorded'}" if reason else f"Outcome: {sample.result_status.replace('_', ' ').title()}", 1.1, 3.62, 9.8, .25, 14, navy)
            return
        heads = ["Chemical", "Notification", "Received", "Reported", "TAT", "STT", "Variance", "Delay reason" if reason else "Outcome"]
        table_height = min(5.35, .31*(len(rows)+1))
        table = slide.shapes.add_table(len(rows)+1, len(heads), Inches(.42), Inches(1.4), Inches(12.45), Inches(table_height)).table
        widths = [2.6,1.25,.85,.85,.55,.55,.7,5.1] if reason else [3.6,1.45,1,1,.65,.65,.85,2.9]
        for i,w in enumerate(widths): table.columns[i].width = Inches(w)
        for c,h in enumerate(heads): table.cell(0,c).text=h; table.cell(0,c).fill.solid(); table.cell(0,c).fill.fore_color.rgb=color(navy)
        for r,s in enumerate(rows,1):
            values=[s.chemical_name, s.notification_no or s.po_number or "—", s.sample_receipt_date.strftime("%d %b") if s.sample_receipt_date else "—", s.report_issue_date.strftime("%d %b") if s.report_issue_date else "—", f"{s.turnaround_days}d", f"{stt(s)}d", f"+{max(s.turnaround_days-stt(s),0)}d", (s.delay_reason or "Not recorded") if reason else s.result_status.replace("_", " ").title()]
            for c,v in enumerate(values):
                table.cell(r,c).text=str(v)
                table.cell(r,c).fill.solid()
                table.cell(r,c).fill.fore_color.rgb = color("F8FBFE" if r % 2 == 0 else "FFFFFF")
        for row in table.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs: p.font.size=Pt(8); p.font.name="Arial"; p.font.color.rgb=color(navy)
        for cell in table.rows[0].cells:
            for p in cell.text_frame.paragraphs: p.font.size=Pt(8); p.font.bold=True; p.font.color.rgb=color("FFFFFF")
    slide = prs.slides.add_slide(blank)
    canvas_background(slide)
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color("EAF4FF")
    rectangle(slide, 12.48, .08, .85, 6.85, "D7EAFB")
    add_text(slide, "QC LABORATORY MONITORING", .76, 1.28, 5.5, .28, 14, blue, True)
    add_text(slide, data["laboratory"]["name"], .76, 1.82, 10.9, .7, 50, navy, True)
    rectangle(slide, .76, 2.75, 1.15, .045, blue)
    add_text(slide, f"Performance Review · {month_start:%B %Y}", .76, 3.08, 8.5, .36, 24, navy, True)
    add_text(slide, "Current quality-control performance, Service Level Agreement compliance and exception review", .76, 3.68, 9.9, .36, 16, grey)
    chrome.add_cover_branding(slide)
    add_text(slide, "Office of Head Corporate Chemistry | Mumbai / Dehradun", .76, 6.45, 7.2, .2, 11, grey)

    slide=prs.slides.add_slide(blank); header(slide, f"{data['laboratory']['name']} | Performance metrics", 2); month=data['month_stt']
    for i,(v,l,t) in enumerate([(data['month_intake'],"Monthly intake",blue),(month['closed'],"Closed reports",navy),(month['within_standard'],"Within STT",green),(month['late'],"Late closures",red),(f"{month['compliance_rate']}%","STT achieved",blue),(f"{month['average_turnaround']} d","Average TAT",navy)]): metric(slide,.6+(i%3)*4.2,1.7+(i//3)*2.0,v,l,t)
    slide=prs.slides.add_slide(blank); header(slide,"Current workload and STT exceptions",3); week=data['week_stt']
    for i,(v,l,t) in enumerate([(data['summary']['total'],"Samples in review",blue),(data['summary']['under_testing'],"Open workload",navy),(week['closed'],"Closed reports",navy),(week['late'],"Closed above STT",red),(f"{week['compliance_rate']}%","Period STT",blue),(len(data['overdue_samples']),"Aged open samples",red)]): metric(slide,.6+(i%3)*4.2,1.7+(i//3)*2.0,v,l,t)
    slide=prs.slides.add_slide(blank); header(slide,"Monthly STT and delay analytics",4)
    metric(slide,.8,1.7,month['within_standard'],"Within applicable STT",green); metric(slide,4.5,1.7,month['late'],"Late closures",red); metric(slide,8.2,1.7,f"{month['compliance_rate']}%","STT achieved",blue)
    reasons={}
    for sample in late: reasons[sample.delay_reason or "No reason recorded"] = reasons.get(sample.delay_reason or "No reason recorded",0)+1
    add_text(slide,"Delay reason distribution",.8,3.45,5.5,.3,18,navy,True)
    for i,(reason,count) in enumerate(sorted(reasons.items(), key=lambda item:(-item[1],item[0]))[:5]): add_text(slide,str(count),.8,3.9+i*.42,.3,.2,16,red if "No reason" in reason else blue,True); add_text(slide,overview_text(reason),1.2,3.9+i*.42,10.9,.25,13,navy)
    slide=prs.slides.add_slide(blank); header(slide,"Exception concentration and review focus",5)
    open_exceptions=data['overdue_samples']; add_text(slide,"Current open samples above the 9-day review threshold",.8,1.55,8.5,.3,18,navy,True)
    if open_exceptions:
        sample=open_exceptions[0]; metric(slide,.8,2.35,len(open_exceptions),"Open STT exceptions",red); add_text(slide,sample.chemical_name,4.5,2.35,4.5,.3,22,navy,True); add_text(slide,f"Notification: {sample.notification_no or sample.po_number or '—'} · Age: {sample.days_open} days",4.5,2.77,5.8,.25,13,grey); add_wrapped_text(slide,f"Reason: {overview_text(sample.delay_reason, 82)}",4.5,3.1,7.7,.5,13,grey)
    else: add_text(slide,"No current open samples are beyond the 9-day threshold.",.8,2.25,7,.3,16,green,True)
    add_text(slide,"Questions for the review",.8,4.35,5,.3,18,navy,True)
    for i,question in enumerate(["Which external testing dependencies require agreed result dates?", "Which late closures need complete delay documentation?", "Who owns each current exception and target closure date?"]): add_text(slide,f"{i+1}",.8,4.8+i*.48,.25,.2,15,blue,True); add_text(slide,question,1.2,4.8+i*.48,9.6,.25,14,navy)
    table_slide("Completed samples above applicable STT", late, 6, True); table_slide("Completed samples within applicable STT", within, 7)
    slide=prs.slides.add_slide(blank); header(slide,"Completed-sample product outcomes",8); passed=sum(s.result_status=="pass" for s in completed); failed=sum(s.result_status=="fail" for s in completed); report_issued=sum(s.result_status=="report_issued" for s in completed)
    metric(slide,.7,1.7,passed,"Pass (product outcome)",green); metric(slide,3.8,1.7,failed,"Fail (product outcome)",red); metric(slide,6.9,1.7,report_issued,"Report issued only",blue); metric(slide,10.0,1.7,len(completed),"Completed samples",navy); add_text(slide,"Product pass/fail is quality context only; report-issued records do not have a recorded pass/fail result.",.8,3.25,11.2,.3,14,grey)
    slide=prs.slides.add_slide(blank); header(slide,"Standard Testing Time (STT) assessment basis",9)
    metric(slide,.8,1.7,month['material_standard_count'],"Material-specific STT",blue); metric(slide,4.5,1.7,month['fallback_stt_count'],"9-day review fallback",blue); metric(slide,8.2,1.7,month['assessed'],"Assessed closures",navy)
    add_text(slide,"Standard Testing Time (STT): the material-specific standard where defined; otherwise the 9-day management-review STT is applied.",.8,3.25,10.8,.35,16,grey)
    output=BytesIO(); prs.save(output); output.seek(0); return output, f"{data['laboratory']['name']} Performance Review {month_start:%b %Y}.pptx"


def build_lab_brief_presentation(lab_code: str, static_folder: str) -> tuple[BytesIO, str]:
    """The local-reporting management brief as a deck, mirroring the page section for section.

    Deliberately narrower than the performance review: this is the one reporting
    week the brief covers, in the same order the page presents it, so a reader
    can follow the screen and the slide interchangeably.
    """
    from pptx import Presentation
    from pptx.util import Inches
    from app.core.services.quality_control import latest_dashboard_data

    data = latest_dashboard_data(lab_code)
    batch = data["batch"]
    if batch is None:
        raise ValueError("Import a local status workbook before downloading the management brief.")

    laboratory, summary, samples = data["laboratory"], data["summary"], data["samples"]
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    chrome = _DeckChrome(
        prs, static_folder,
        f"Source: {laboratory['name']} local status workbook \u00b7 {batch.report_label}",
    )
    navy, blue, red, green, grey = chrome.NAVY, chrome.BLUE, chrome.RED, chrome.GREEN, chrome.GREY

    # 01 · Cover
    slide = chrome.new_slide("Quality Control Laboratory — Management Brief", 1)
    chrome.add_text(slide, laboratory["name"], .42, 1.85, 9.0, .6, 34, navy, True)
    chrome.add_text(slide, laboratory["location"], .42, 2.55, 9.0, .35, 16, grey)
    chrome.add_text(slide, batch.report_label, .42, 3.05, 11.0, .4, 20, blue, True)
    chrome.rectangle(slide, .42, 3.62, 3.2, .05, blue)

    # 02 · The four figures the brief leads with
    slide = chrome.new_slide("Current reporting position", 2)
    turnaround = f"{summary['average_turnaround']} d" if summary["average_turnaround"] is not None else "—"
    chrome.metric(slide, .8, 1.9, summary["total"], "Reported sample load", blue)
    chrome.metric(slide, 4.3, 1.9, summary["under_testing"], "Open workload", red if summary["under_testing"] else navy)
    chrome.metric(slide, 7.8, 1.9, f"{summary['passed']}/{summary['failed']}", "Pass / fail reports issued", green)
    chrome.metric(slide, .8, 3.9, turnaround, "Average turnaround, issued", navy)
    chrome.metric(slide, 4.3, 3.9, summary["delayed_open"], "Aged beyond target", red if summary["delayed_open"] else green)
    chrome.metric(slide, 7.8, 3.9, summary["issued"], "Reports issued in this period", navy)

    # 03 · Management attention required
    attention = [s for s in samples if s.result_status == "under_testing" and s.delay_reason]
    slide = chrome.new_slide("Management attention required", 3)
    if attention:
        y = 1.75
        for sample in attention[:8]:
            received = sample.sample_receipt_date.strftime("%d %b %Y") if sample.sample_receipt_date else "date not recorded"
            chrome.rectangle(slide, .8, y, 11.6, .04, red)
            chrome.add_text(slide, sample.chemical_name, .8, y + .12, 5.4, .3, 16, navy, True)
            chrome.add_text(slide, f"Received {received}", 6.4, y + .14, 3.0, .25, 12, grey)
            chrome.add_wrapped_text(slide, " ".join(str(sample.delay_reason).split()), .8, y + .48, 11.5, .42, 12, grey)
            y += 1.06
        if len(attention) > 8:
            chrome.add_text(slide, f"+ {len(attention) - 8} further open samples with recorded remarks", .8, y + .05, 8.0, .3, 12, grey)
    else:
        chrome.add_text(slide, "No open sample has a recorded delay remark in this reporting period.", .8, 2.0, 10.5, .35, 18, green, True)

    # 04 · Issued reports
    issued = [s for s in samples if s.result_status != "under_testing"]
    slide = chrome.new_slide("Issued reports", 4)
    if issued:
        rows = [["Material", "Outcome", "Report date", "Turnaround"]] + [
            [
                s.chemical_name,
                "Report issued" if s.result_status == "report_issued" else s.result_status.title(),
                s.report_issue_date.strftime("%d %b %Y") if s.report_issue_date else "—",
                f"{s.turnaround_days} days" if s.turnaround_days is not None else "—",
            ]
            for s in issued[:14]
        ]
        table = slide.shapes.add_table(
            len(rows), 4, Inches(.8), Inches(1.7), Inches(11.6), Inches(.34 * len(rows)),
        ).table
        for width, index in zip((5.0, 2.4, 2.2, 2.0), range(4)):
            table.columns[index].width = Inches(width)
        for r, row in enumerate(rows):
            for c, value in enumerate(row):
                cell = table.cell(r, c)
                cell.text = str(value)
                cell.fill.solid()
                cell.fill.fore_color.rgb = chrome.color(navy if r == 0 else ("FFFFFF" if r % 2 else "F2F4F7"))
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = chrome._Pt(12 if r else 11)
                paragraph.font.bold = r == 0
                paragraph.font.name = "Arial"
                paragraph.font.color.rgb = chrome.color("FFFFFF" if r == 0 else navy)
        if len(issued) > 14:
            chrome.add_text(slide, f"+ {len(issued) - 14} further issued reports in the full register", .8, 1.75 + .34 * len(rows), 8.0, .3, 12, grey)
    else:
        chrome.add_text(slide, "No reports were issued in this reporting period.", .8, 2.0, 10.5, .35, 18, grey, True)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output, f"{laboratory['name']} Management Brief {batch.week_end:%d %b %Y}.pptx"


def build_portfolio_management_presentation(static_folder: str, reporting_week_end=None, lab_codes: set[str] | None = None) -> tuple[BytesIO, str]:
    """Create an on-demand, consolidated management-review presentation.

    ``lab_codes`` builds it for the chosen laboratories instead of the whole
    network. Every figure then counts only those laboratories, and the deck says
    so on its cover, on every slide's source line, and in its filename.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
    from app.core.services.quality_control import portfolio_management_data

    data = portfolio_management_data(reporting_week_end, lab_codes)
    if not data["reporting_labs"]:
        raise ValueError(
            "None of the selected laboratories submitted a workbook for this reporting period."
            if lab_codes is not None else
            "Import a local status workbook for at least one laboratory before downloading a presentation."
        )
    names = [review["laboratory"]["name"] for review in data["laboratory_reviews"]]
    if lab_codes is None:
        scope_heading, scope_note, scope_stem = "ALL ONGC", "the whole laboratory network", "QC Portfolio"
    elif len(names) == 1:
        scope_heading = scope_note = scope_stem = names[0]
    else:
        listed = ", ".join(names[:4]) + (f" and {len(names) - 4} more" if len(names) > 4 else "")
        scope_heading, scope_stem = f"{len(names)} LABORATORIES", f"{len(names)} laboratories"
        scope_note = f"{len(names)} laboratories — {listed}"

    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5); blank = prs.slide_layouts[6]
    navy, blue, red, green, grey, border = "071D42", "1976D2", "C53B3B", "18794E", "526173", "D7E2EE"
    def color(value): return RGBColor.from_string(value)
    def add_text(slide, value, x, y, w, h, size=18, tone=navy, bold=False, wrap=False):
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); shape.text_frame.word_wrap = wrap
        paragraph = shape.text_frame.paragraphs[0]; paragraph.text = str(value); paragraph.font.size = Pt(size); paragraph.font.bold = bold; paragraph.font.color.rgb = color(tone); paragraph.font.name = "Arial"
        return shape
    def rectangle(slide, x, y, w, h, fill, line=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)); shape.fill.solid(); shape.fill.fore_color.rgb = color(fill); shape.line.color.rgb = color(line or fill)
        return shape
    def concise(value, limit=55):
        text = " ".join(str(value or "Not recorded").split())
        return text if len(text) <= limit else f"{text[:limit + 1].rsplit(' ', 1)[0]}…"
    corporate_chemistry_logo = Path(static_folder) / "images" / "ongc-corporate-chemistry-logo.png"
    ongc_logo = Path(static_folder) / "images" / "ongc-official-logo.png"
    period = data["reporting_period"]
    period_label = f"{period['start']:%d %b} – {period['end']:%d %b %Y}" if period else "Latest reporting period"
    def background(slide):
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color("FFFFFF")
        rectangle(slide, 0, 0, 13.333, .08, navy); rectangle(slide, .42, 1.26, 12.45, .015, border); rectangle(slide, .42, 6.93, 12.45, .015, border)
    def header(slide, title, page):
        background(slide); add_text(slide, "ONGC CORPORATE CHEMISTRY · QC LABORATORY MONITORING", .42, .27, 7.6, .24, 11, blue, True); add_text(slide, title, .42, .7, 11.35, .46, 27, navy, True)
        if ongc_logo.exists(): slide.shapes.add_picture(str(ongc_logo), Inches(11.15), Inches(.21), width=Inches(.9), height=Inches(.51))
        if corporate_chemistry_logo.exists(): slide.shapes.add_picture(str(corporate_chemistry_logo), Inches(12.24), Inches(.16), width=Inches(.55), height=Inches(.55))
        add_text(slide, f"Source: Selected reporting period · {period_label} · Scope: {scope_note}", .42, 7.08, 9.6, .16, 8, grey); add_text(slide, f"{page:02d}", 12.5, 7.08, .25, .16, 8, grey)
    def metric(slide, x, y, value, label, tone=navy):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x-.18), Inches(y-.18), Inches(3.45), Inches(1.35)); card.fill.solid(); card.fill.fore_color.rgb = color("EAF4FF" if tone == blue else "FCEBEC" if tone == red else "EAF7F0" if tone == green else "F2F4F7"); card.line.color.rgb = color(border)
        add_text(slide, value, x, y, 2.8, .4, 27, tone, True); add_text(slide, label, x, y+.52, 3.0, .25, 13, navy, True)
    def table(slide, headers, rows, widths, y=1.45, font_size=9):
        if not rows:
            add_text(slide, "No records require management attention for this view.", .8, 2.1, 10.5, .35, 18, green, True); return
        height = min(5.25, .34 * (len(rows) + 1)); table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(.42), Inches(y), Inches(12.45), Inches(height)).table
        for index, width in enumerate(widths): table_shape.columns[index].width = Inches(width)
        for col, label in enumerate(headers):
            cell = table_shape.cell(0, col); cell.text = label; cell.fill.solid(); cell.fill.fore_color.rgb = color(navy)
        for row_index, values in enumerate(rows, 1):
            for col, value in enumerate(values):
                cell = table_shape.cell(row_index, col); cell.text = str(value); cell.fill.solid(); cell.fill.fore_color.rgb = color("F8FBFE" if row_index % 2 == 0 else "FFFFFF")
        for row in table_shape.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(font_size); paragraph.font.name = "Arial"; paragraph.font.color.rgb = color(navy)
        for cell in table_shape.rows[0].cells:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True; paragraph.font.color.rgb = color("FFFFFF")

    def paginated_table(title, headers, rows, widths, font_size=8, rows_per_slide=12):
        """Add every register row, continuing the exception table across slides."""
        if not rows:
            slide = prs.slides.add_slide(blank)
            header(slide, title, len(prs.slides))
            table(slide, headers, [], widths, font_size=font_size)
            return
        total = len(rows)
        for start in range(0, total, rows_per_slide):
            end = min(start + rows_per_slide, total)
            slide = prs.slides.add_slide(blank)
            suffix = f" ({start + 1}–{end} of {total})" if total > rows_per_slide else ""
            header(slide, f"{title}{suffix}", len(prs.slides))
            table(slide, headers, rows[start:end], widths, font_size=font_size)

    slide = prs.slides.add_slide(blank); background(slide); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color("EAF4FF"); rectangle(slide, 12.48, .08, .85, 6.85, "D7EAFB")
    add_text(slide, f"QC LABORATORY MONITORING · {scope_heading}", .76, 1.28, 8.5, .28, 14, blue, True); add_text(slide, "Portfolio Management Review", .76, 1.82, 10.5, .7, 50, navy, True); rectangle(slide, .76, 2.75, 1.15, .045, blue)
    add_text(slide, period_label, .76, 3.08, 8.5, .36, 24, navy, True); add_text(slide, f"Consolidated position across {data['reporting_labs']} reporting laborator{'y' if data['reporting_labs'] == 1 else 'ies'} — {scope_note}", .76, 3.68, 10.5, .36, 16, grey, wrap=True)
    rectangle(slide, 10.0, .46, 2.32, 1.2, "FFFFFF", border)
    if ongc_logo.exists(): slide.shapes.add_picture(str(ongc_logo), Inches(10.16), Inches(.68), width=Inches(.9), height=Inches(.51))
    if corporate_chemistry_logo.exists(): slide.shapes.add_picture(str(corporate_chemistry_logo), Inches(11.16), Inches(.53), width=Inches(1.05), height=Inches(1.05))
    add_text(slide, "Office of Head Corporate Chemistry | Mumbai / Dehradun", .76, 6.45, 7.2, .2, 11, grey)

    summary = data["summary"]
    slide = prs.slides.add_slide(blank); header(slide, "Portfolio delivery position", 2)
    for index, item in enumerate([(summary["total"], "Samples in review", blue), (summary["under_testing"], "Open workload", navy), (summary["delayed_open"], "Aged open > 9 days", red), (f"{summary['passed']} / {summary['failed']}", "Pass / fail reports", green), (f"{summary['average_turnaround']} d" if summary["average_turnaround"] is not None else "—", "Average turnaround", blue), (data["reporting_labs"], "Laboratories submitted", navy)]):
        metric(slide, .6 + (index % 3) * 4.2, 1.7 + (index // 3) * 2.0, item[0], item[1], item[2])
    if data["missing_submissions"]:
        missing_names = ", ".join(review["laboratory"]["name"] for review in data["missing_submissions"])
        add_text(slide, f"Excluded from this reporting period (no submission): {missing_names}", .8, 5.0, 11.2, .3, 13, red, True)
    else:
        add_text(slide, "All configured laboratories submitted data for the selected reporting period.", .8, 5.0, 11.2, .3, 13, green, True)

    slide = prs.slides.add_slide(blank); header(slide, "Laboratory performance at a glance", 3)
    rows = []
    for review in data["laboratory_reviews"]:
        batch, item = review["batch"], review["summary"]
        status = "Submitted" if batch else (f"Not received · last {review['latest_available_batch'].week_end:%d %b}" if review["latest_available_batch"] else "No submission")
        rows.append([review["laboratory"]["name"], status, item["total"] if batch else "—", item["under_testing"] if batch else "—", item["delayed_open"] if batch else "—", f"{item['passed']} / {item['failed']}" if batch else "—", f"{item['average_turnaround']} d" if batch and item["average_turnaround"] is not None else "—"])
    table(slide, ["Laboratory", "Submission status", "Samples", "Open", "Aged >9d", "Pass / fail", "Average TAT"], rows, [2.45, 2.1, .85, .8, .95, 1.15, 1.1], font_size=9)

    standard_rows = [item for item in data["completed_testing"] if item["standard_days"] is not None and item["variance_days"] is not None]
    late_rows = [item for item in standard_rows if item["variance_days"] > 0]
    within = sum(item["variance_days"] <= 0 for item in standard_rows)
    compliance = round(within / len(standard_rows) * 100, 1) if standard_rows else None
    slide = prs.slides.add_slide(blank); header(slide, "Service Level Agreement performance", 4)
    for index, item in enumerate([(len(standard_rows), "Completed tests assessed", navy), (within, "Within approved standard", green), (len(late_rows), "Above approved standard", red), (f"{compliance}%" if compliance is not None else "—", "Service Level Agreement achieved", blue)]):
        metric(slide, .7 + index * 3.1, 1.7, item[0], item[1], item[2])
    add_text(slide, "A material-specific approved standard is used where available; no fallback is included in this cross-laboratory comparison.", .8, 3.35, 11.1, .3, 14, grey)

    late_table_rows = [[data["laboratories_by_code"].get(item["sample"].lab_code, {"name": item["sample"].lab_code})["name"], item["sample"].chemical_name, item["sample"].notification_no or item["sample"].po_number or "—", f"{item['sample'].turnaround_days} d", f"{item['standard_days']} d", f"+{item['variance_days']} d", concise(item["sample"].delay_reason)] for item in late_rows]
    paginated_table("Completed tests above STT", ["Laboratory", "Chemical", "Notification", "Actual", "Standard", "Variance", "Delay reason"], late_table_rows, [1.75, 2.15, 1.2, .75, .85, .85, 4.9])

    open_rows = [[data["laboratories_by_code"].get(sample.lab_code, {"name": sample.lab_code})["name"], sample.chemical_name, sample.notification_no or sample.po_number or "—", sample.sample_receipt_date.strftime("%d %b %Y") if sample.sample_receipt_date else "—", f"{sample.days_open} d", concise(sample.delay_reason)] for sample in data["overdue_samples"]]
    paginated_table("Open samples above 9-day threshold", ["Laboratory", "Material", "Notification", "Received", "Age", "Remarks"], open_rows, [1.85, 2.4, 1.35, 1.15, .7, 4.85])

    chemical_rows = [[item["chemical_name"], ", ".join(item["laboratories"]), item["total"], item["passed"], item["failed"], item["under_testing"], f"{item['average_actual']} d" if item["average_actual"] is not None else "—", f"{item['standard_days']} d" if item["standard_days"] is not None else "—"] for item in data["weekly_chemical_metrics"]]
    paginated_table("Reporting-period chemical performance", ["Chemical", "Laboratories", "Samples", "Pass", "Fail", "Open", "Average time", "Standard"], chemical_rows, [2.35, 3.35, .75, .65, .65, .65, 1.15, 1.15])

    output = BytesIO(); prs.save(output); output.seek(0)
    return output, f"{scope_stem} Management Review {period['end']:%b %Y}.pptx"


def build_sap_portfolio_management_presentation(
    static_folder: str, lab_codes: set[str] | None = None,
) -> tuple[BytesIO, str]:
    """Create the senior-management deck from the current SAP snapshots only."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from app.core.services.sap_quality_control import (
        non_sap_register_data, sap_management_data,
    )

    data = sap_management_data(lab_codes)
    # The declared non-SAP register is read separately and stays separate: it
    # never enters the SAP counts, but it is the rest of the bench's load and
    # the management deck is now the only deck that carries it.
    non_sap = non_sap_register_data(lab_codes)
    if not data["reporting_labs"]:
        raise ValueError("Import paired SAP exports for at least one laboratory before downloading the management presentation.")

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    scope_labs = data["scope_laboratories"]
    scope_label = "All SAP laboratories" if lab_codes is None else ", ".join(lab["name"] for lab in scope_labs)
    chrome = _DeckChrome(
        prs, static_folder,
        f"Source: latest paired SAP Inspection Lots and Notifications exports · {data['source_as_of_label']}",
    )
    navy, blue, red, green, grey = chrome.NAVY, chrome.BLUE, chrome.RED, chrome.GREEN, chrome.GREY

    def concise(value, limit=36):
        text = " ".join(str(value or "—").split())
        return text if len(text) <= limit else f"{text[:limit].rsplit(' ', 1)[0]}…"

    def table(slide, headers, rows, widths, *, y=1.5, font_size=9):
        if not rows:
            chrome.add_text(slide, "No SAP records are available for this view.", .8, 2.1, 10.5, .35, 18, green, True)
            return
        height = min(5.15, .34 * (len(rows) + 1))
        shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(.42), Inches(y), Inches(12.45), Inches(height))
        table_shape = shape.table
        for index, width in enumerate(widths):
            table_shape.columns[index].width = Inches(width)
        for column, label in enumerate(headers):
            cell = table_shape.cell(0, column)
            cell.text = label
            cell.fill.solid()
            cell.fill.fore_color.rgb = chrome.color(navy)
        for row_index, values in enumerate(rows, 1):
            for column, value in enumerate(values):
                cell = table_shape.cell(row_index, column)
                cell.text = str(value)
                cell.fill.solid()
                cell.fill.fore_color.rgb = chrome.color("F8FBFE" if row_index % 2 == 0 else "FFFFFF")
        for row in table_shape.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(font_size)
                    paragraph.font.name = "Arial"
                    paragraph.font.color.rgb = chrome.color(navy)
        for cell in table_shape.rows[0].cells:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.color.rgb = chrome.color("FFFFFF")

    # 01 · Cover
    slide = prs.slides.add_slide(chrome.blank)
    chrome.canvas_background(slide)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = chrome.color("EAF4FF")
    chrome.rectangle(slide, 12.48, .08, .85, 6.85, "D7EAFB")
    chrome.add_text(slide, "QC LABORATORY MONITORING · SAP QM", .76, 1.28, 8.6, .28, 14, blue, True)
    chrome.add_text(slide, "Management Review", .76, 1.82, 10.5, .7, 50, navy, True)
    chrome.rectangle(slide, .76, 2.75, 1.15, .045, blue)
    chrome.add_text(slide, data["source_as_of_label"], .76, 3.08, 8.6, .36, 24, navy, True)
    chrome.add_wrapped_text(slide, f"Scope: {scope_label}. All figures and action registers are derived from the latest paired SAP exports; local workbook data is excluded.", .76, 3.68, 8.7, .7, 16, grey)
    chrome.add_cover_branding(slide)
    chrome.add_text(slide, "Office of Head Corporate Chemistry | Mumbai / Dehradun", .76, 6.45, 7.2, .2, 11, grey)

    # 02 · Position
    kpis = data["kpis"]
    slide = chrome.new_slide("Official SAP position", 2)
    cards = [
        (kpis["total"], "SAP monitoring records", blue),
        (kpis["actionable_open"], "Actionable SAP-open", red if kpis["actionable_open"] else green),
        (kpis["stt_overdue"], "Past STT", red if kpis["stt_overdue"] else green),
        (kpis["awaiting_lab"], "Awaiting lab follow-up", blue),
        (f"{kpis['accepted']} / {kpis['rejected']}", "UD A / UD R", green),
        (kpis["completed"], "Officially complete", green),
    ]
    for index, (value, label, tone) in enumerate(cards):
        chrome.metric(slide, .7 + (index % 3) * 4.2, 1.65 + (index // 3) * 2.0, value, label, tone)
    chrome.add_text(slide, f"Snapshot coverage: {data['reporting_labs']} of {data['configured_labs']} configured SAP laboratories. QC-admin exclusions: {kpis['excluded']}; exclusions requiring renewed review: {kpis['exclusion_review']}.", .75, 5.9, 11.6, .3, 13, grey)

    # 03 · Laboratory overview
    slide = chrome.new_slide("Laboratory SAP snapshot coverage", 3)
    lab_rows = []
    for review in data["laboratory_reviews"]:
        if review["batch"] is None:
            lab_rows.append([review["laboratory"]["name"], "—", "Awaiting snapshot", "—", "—", "—", "—"])
            continue
        item = review["kpis"]
        lab_rows.append([
            review["laboratory"]["name"], review["batch"].plant_code,
            review["batch"].as_of_date.strftime("%d %b %Y"), item["total"],
            item["actionable_open"], item["stt_overdue"], item["awaiting_lab"],
        ])
    table(slide, ["Laboratory", "Plant", "SAP as of", "Records", "Actionable open", "Past STT", "Awaiting lab"], lab_rows, [2.8, 1.0, 1.45, 1.15, 1.65, 1.25, 1.15], y=1.55, font_size=9)

    # 04 · Work-centre exposure
    slide = chrome.new_slide("Open workload by SAP work center", 4)
    centre_rows = [[
        item["name"], ", ".join(item["laboratories"]), item["open"],
        item["stt_overdue"], item["awaiting_lab"],
    ] for item in data["work_centers"][:14]]
    table(slide, ["SAP work center", "Laboratories", "Open", "Past STT", "No lab update"], centre_rows, [3.7, 3.65, 1.35, 1.65, 2.1], y=1.55, font_size=10)

    # 05+ · Complete current action register, grouped for each laboratory's
    # hand-off by Corporate Specification sub-group rather than as one mixed
    # cross-laboratory queue.
    action_groups = _sap_presentation_action_groups(data)
    page_index = 5
    if not action_groups:
        slide = chrome.new_slide("All actionable SAP-open items", page_index)
        table(slide, ["Inspection lot", "Notification", "Material", "Specification", "Work center", "STT due", "Lab follow-up"], [], [1.35, 1.35, 2.25, 1.8, 1.65, 1.45, 2.6], y=1.48, font_size=8)
        page_index += 1
    for group in action_groups:
        total = len(group["entries"])
        for start, entries in enumerate(_paginated_rows(group["entries"], 10)):
            first = start * 10 + 1
            last = first + len(entries) - 1
            suffix = f" ({first}–{last} of {total})" if total > 10 else ""
            title = f"{group['laboratory']['name']} · {group['subgroup_label']}{suffix}"
            slide = chrome.new_slide(title, page_index)
            page_index += 1
            chrome.add_text(
                slide,
                "Actionable SAP-open samples — grouped by laboratory and Corporate Specification sub-group.",
                .45, 1.28, 10.6, .15, 8, grey,
            )
            rows = []
            for item in entries:
                record, update = item["record"], item["lab_update"]
                if item["stt_due_date"]:
                    stt_due = item["stt_due_date"].strftime("%d %b %Y")
                    if item["stt_overdue"]:
                        stt_due += f" · {item['stt_variance_days']} d over"
                elif item["stt_days"] is not None:
                    stt_due = f"STT {item['stt_days']} d · no start"
                else:
                    stt_due = "STT not defined"
                follow_up = item["reconciliation_label"]
                if update and update.expected_completion_date:
                    follow_up += f" · ETA {update.expected_completion_date:%d %b}"
                rows.append([
                    record.inspection_lot_number or "—", record.notification_no or "—",
                    concise(record.material_description, 31),
                    concise(item["specification_no"] or "Not in Corporate Specification", 28),
                    concise(record.work_center or "Not assigned", 24), stt_due, concise(follow_up, 34),
                ])
            table(slide, ["Inspection lot", "Notification", "Material", "Specification", "Work center", "STT due", "Lab follow-up"], rows, [1.35, 1.35, 2.25, 1.8, 1.65, 1.45, 2.6], y=1.48, font_size=8)

    # Non-SAP register · declared samples with no SAP record, kept apart from
    # every count above.
    non_sap_kpis = non_sap["non_sap_kpis"]
    if non_sap_kpis["total"]:
        slide = chrome.new_slide("Non-SAP samples — separate declared register", page_index)
        page_index += 1
        cards = [
            (non_sap_kpis["total"], "Declared non-SAP samples", blue),
            (non_sap_kpis["pending"], "Still with the laboratory", red if non_sap_kpis["pending"] else green),
            (non_sap_kpis["overdue"], "Past the declared ETA", red if non_sap_kpis["overdue"] else green),
        ]
        for index, (value, label, tone) in enumerate(cards):
            chrome.metric(slide, .8 + index * 4.2, 1.62, value, label, tone)
        # A declared result, never an SAP usage decision, so it is worded and
        # counted apart from the UD figures earlier in the deck.
        chrome.add_text(slide, f"Declared results to date: {non_sap_kpis['closed_pass']} pass · {non_sap_kpis['closed_fail']} fail.", .75, 3.15, 11.6, .28, 14, navy, True)
        lab_rows = [[
            item["laboratory"]["name"], item["total"], item["pending"],
            item["overdue"], item["closed_pass"], item["closed_fail"],
        ] for item in non_sap["non_sap_by_laboratory"]]
        table(slide, ["Laboratory", "Declared", "Pending", "Past ETA", "Declared pass", "Declared fail"], lab_rows, [3.5, 1.7, 1.7, 1.7, 2.0, 1.85], y=3.6, font_size=10)
        chrome.add_text(slide, "These samples are not in SAP and are excluded from every SAP figure in this deck.", .75, 6.65, 11.6, .25, 11, grey)

        pending_entries = [item for item in non_sap["non_sap_entries"] if not item["is_closed"]]
        total_pending = len(pending_entries)
        for start, entries in enumerate(_paginated_rows(pending_entries, 10)):
            first = start * 10 + 1
            last = first + len(entries) - 1
            suffix = f" ({first}–{last} of {total_pending})" if total_pending > 10 else ""
            slide = chrome.new_slide(f"Non-SAP samples awaiting return{suffix}", page_index)
            page_index += 1
            rows = []
            for item in entries:
                sample = item["sample"]
                eta = sample.expected_completion_date.strftime("%d %b %Y") if sample.expected_completion_date else "ETA not declared"
                if item["is_overdue"]:
                    eta += " · overdue"
                rows.append([
                    item["laboratory"]["name"],
                    concise(sample.sample_reference, 18),
                    concise(sample.chemical_name, 30),
                    concise(item["status_label"], 26),
                    eta,
                    concise(sample.action_owner or sample.delay_reason or "—", 26),
                ])
            table(slide, ["Laboratory", "Local reference", "Material / sample", "Declared stage", "Expected completion", "Owner / constraint"], rows, [2.35, 1.85, 2.75, 2.15, 1.9, 1.45], y=1.55, font_size=9)
            chrome.add_text(slide, "Laboratories record these returns on their own dashboard; this page reports what is outstanding.", .65, 6.5, 11.6, .25, 10, grey)

    slide = chrome.new_slide("SAP usage decisions and daily movement", page_index)
    for index, item in enumerate(data["usage_decisions"]):
        tone = green if item["tone"] == "success" else red if item["tone"] == "danger" else navy
        chrome.metric(slide, .8 + index * 3.7, 1.7, item["count"], item["label"], tone)
    movement_rows = [[
        item["laboratory"]["name"], item["batch"].as_of_date.strftime("%d %b"),
        item["current_open"], item["previous_open"] if item["previous_open"] is not None else "—",
        f"{item['open_change']:+d}" if item["open_change"] is not None else "First snapshot",
    ] for item in data["trend"]]
    chrome.add_text(slide, "Change from the previous SAP snapshot", .8, 3.55, 5.6, .3, 18, navy, True)
    table(slide, ["Laboratory", "Current", "SAP-open", "Previous open", "Change"], movement_rows, [3.4, 2.0, 2.0, 2.35, 2.7], y=4.0, font_size=10)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    filename_date = data["source_dates"][0].strftime("%d %b %Y") if len(data["source_dates"]) == 1 else "Latest"
    if lab_codes is None:
        filename_scope = "QC SAP"
    elif len(scope_labs) == 1:
        filename_scope = scope_labs[0]["name"]
    else:
        filename_scope = f"{len(scope_labs)} SAP Laboratories"
    return output, f"{filename_scope} Management Review {filename_date}.pptx"
